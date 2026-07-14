"""Tools the ingestion agent can call.

Each tool has:
- A Claude tool schema (for the tool-use API)
- An async executor that runs the actual work

All tools share a single `job_id` and read/write via `services.job_cache`,
which is disk-backed and therefore survives backend restarts. Tools delegate
heavy lifting to the existing stage modules so this is an orchestration
layer, not a reimplementation.
"""
from __future__ import annotations

import asyncio
import logging
import time
from pathlib import Path
from typing import Any, Awaitable, Callable

import services.job_cache as cache

log = logging.getLogger(__name__)

# Cap on how many pages we'll inspect for the OCR-fraction signal. For docs
# larger than this we sample evenly (first 20 + last 10 + spaced middles) so
# inspection stays under ~30s even on a 1000-page brick.
_OCR_SCAN_PAGE_CAP    = 200
_OCR_TEXT_THRESHOLD   = 30   # chars/page below this == "needs OCR"
_OCR_FRACTION_THRESH  = 0.20 # >20% pages-needing-OCR => use Docling


def _scan_pdf_ocr_signal(pages) -> dict:
    """Return ocr_fraction + diagnostics by checking extractable text per page.

    For docs up to `_OCR_SCAN_PAGE_CAP` pages we scan every page. For larger
    docs we sample: first 20 pages + last 10 pages + spaced middles, up to
    the cap. A page is "low-text" if pdfplumber extracts < `_OCR_TEXT_THRESHOLD`
    chars — strong proxy for "this page is scanned image, not born-digital".

    Returns:
      ocr_fraction       — float in [0, 1]: pages_low_text / pages_sampled
      pages_needing_ocr  — int: pages flagged
      pages_sampled      — int: number of pages we actually checked
      sampled_indices    — list[int]: page numbers (1-indexed) we sampled
    """
    total = len(pages)
    if total == 0:
        return {"ocr_fraction": 0.0, "pages_needing_ocr": 0, "pages_sampled": 0, "sampled_indices": []}

    if total <= _OCR_SCAN_PAGE_CAP:
        sample_indices = list(range(total))
    else:
        head_n, tail_n = 20, 10
        middle_budget = _OCR_SCAN_PAGE_CAP - head_n - tail_n
        head = list(range(head_n))
        tail = list(range(total - tail_n, total))
        middle_range = range(head_n, total - tail_n)
        if middle_range:
            step = max(1, len(middle_range) // middle_budget)
            middle = list(middle_range)[::step][:middle_budget]
        else:
            middle = []
        sample_indices = sorted(set(head + middle + tail))

    low_text = 0
    for idx in sample_indices:
        try:
            text = (pages[idx].extract_text() or "").strip()
        except Exception:
            text = ""
        if len(text) < _OCR_TEXT_THRESHOLD:
            low_text += 1

    sampled = len(sample_indices)
    return {
        "ocr_fraction":      round(low_text / sampled, 3) if sampled else 0.0,
        "pages_needing_ocr": low_text,
        "pages_sampled":     sampled,
        "sampled_indices":   [i + 1 for i in sample_indices],  # 1-indexed for humans
    }

# ── Tool schemas (passed to Anthropic tool-use API) ───────────────────────────

TOOL_SCHEMAS: list[dict] = [
    {
        "name": "inspect_document",
        "description": (
            "ALWAYS CALL FIRST. Returns the document's format, size, page count, "
            "whether text is natively extractable, whether it contains images, "
            "whether it appears scanned (low extractable text per page), and a "
            "short text sample. Use this signal to pick the right parser."
        ),
        "input_schema": {
            "type": "object",
            "properties": {},
            "required": [],
        },
    },
    {
        "name": "parse_pdf_native",
        "description": (
            "Fast text+table extraction for born-digital PDFs using pdfplumber. "
            "Use when inspect_document reports has_extractable_text=true and "
            "is_scanned=false. Does NOT do OCR — will miss content on scanned pages."
        ),
        "input_schema": {
            "type": "object",
            "properties": {},
            "required": [],
        },
    },
    {
        "name": "parse_with_docling",
        "description": (
            "Heavy ML-based parsing using IBM Docling: TableFormer for tables, "
            "RapidOCR for scanned pages, layout-aware text extraction. SLOW: "
            "~7s per page on CPU. Use only when the doc is scanned OR has "
            "complex tables that pdfplumber will mangle. Returns chunks, tables, "
            "and OCR'd content."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "do_ocr": {
                    "type": "boolean",
                    "description": "Run OCR. Set true for scanned PDFs.",
                    "default": True,
                },
            },
            "required": [],
        },
    },
    {
        "name": "parse_office_document",
        "description": (
            "Parse a DOCX, PPTX, XLSX, or HTML file using format-specific "
            "libraries (python-docx, openpyxl, etc.). Use for non-PDF office "
            "formats."
        ),
        "input_schema": {
            "type": "object",
            "properties": {},
            "required": [],
        },
    },
    {
        "name": "parse_with_vision_ocr",
        "description": (
            "Last-resort OCR via Claude vision: renders every PDF page to a PNG "
            "and asks Claude to transcribe text + extract tables. Use when "
            "inspect_document reports the document is fully scanned, has "
            "CID-encoded/broken fonts (sample_text full of '(cid:N)' codes), or "
            "when parse_with_docling has already timed out. Costs ~$0.01-0.02 "
            "per page (Haiku vision) — cap is 200 pages. Slower per-page than "
            "RapidOCR but produces dramatically better English text and "
            "recovers table structure as markdown. Sets up text_blocks + tables "
            "exactly like the other parsers, so chunk_text/embed_and_index "
            "work without changes after this."
        ),
        "input_schema": {
            "type": "object",
            "properties": {},
            "required": [],
        },
    },
    {
        "name": "chunk_text",
        "description": (
            "Smart heading-aware chunking (~300 tokens, 50-token overlap). "
            "Call AFTER a parser has populated text_blocks. Required before "
            "embed_and_index."
        ),
        "input_schema": {
            "type": "object",
            "properties": {},
            "required": [],
        },
    },
    {
        "name": "describe_tables",
        "description": (
            "MANDATORY when inspect_document reported has_tables=true. "
            "Generates one-sentence Claude descriptions for every structured "
            "table the parser extracted, then builds a `table_summary` chunk "
            "per table (description + columns + sample rows) and pushes them "
            "to the embedding queue. This is how table semantics become "
            "searchable in the vector DB — without this, only raw row text "
            "is embedded and questions like \"what does the revenue table show\" "
            "can't be answered from the chunks. Batched Claude call — cheap."
        ),
        "input_schema": {"type": "object", "properties": {}, "required": []},
    },
    {
        "name": "caption_images",
        "description": (
            "Process visual content ONLY (does NOT touch tables — use "
            "describe_tables for that). Two things: (1) one-sentence Claude "
            "vision caption per embedded image, (2) structured OCR on any "
            "image-only/scanned pages (extracting prose + any tables visible). "
            "Skip if has_images=false. Skip if images are purely decorative "
            "(logos, page borders) — but use for charts, diagrams, infographics, "
            "scanned pages. Bounded 4-concurrent vision calls, per-image retry."
        ),
        "input_schema": {"type": "object", "properties": {}, "required": []},
    },
    {
        "name": "embed_and_index",
        "description": (
            "Embed chunks with text-embedding-3-large and upsert into Qdrant. "
            "This is the terminal step for vector retrieval — required for the "
            "doc to be queryable. Call after chunk_text."
        ),
        "input_schema": {
            "type": "object",
            "properties": {},
            "required": [],
        },
    },
    {
        "name": "store_tables_sql",
        "description": (
            "Index extracted structured tables into a per-job SQLite database "
            "so SQL-routed questions ('what was the Q3 ARR?') can be answered "
            "exactly. Call only if a parser found structured tables with "
            "headers and rows."
        ),
        "input_schema": {
            "type": "object",
            "properties": {},
            "required": [],
        },
    },
    {
        "name": "extract_entities",
        "description": (
            "Build a knowledge graph: named entities (PERSON, ORG, GPE, DATE, "
            "MONEY) and their co-occurrences across chunks. Useful for "
            "entity-anchored questions. Skip for table-only documents or "
            "purely procedural text."
        ),
        "input_schema": {
            "type": "object",
            "properties": {},
            "required": [],
        },
    },
    {
        "name": "finalize",
        "description": (
            "Call when ingestion is complete. Provide a short summary of what "
            "you did and any caveats (e.g., 'images skipped because purely "
            "decorative'). This stops the agent loop."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "summary": {
                    "type": "string",
                    "description": "1-2 sentence summary of what was ingested and what was skipped.",
                },
            },
            "required": ["summary"],
        },
    },
]


# ── Tool executors ─────────────────────────────────────────────────────────────


async def _inspect_document(filepath: Path, job_id: str, **_) -> dict:
    """Inspect a doc and report what it looks like — drives the agent's decisions."""
    ext = filepath.suffix.lower().lstrip(".")
    size_bytes = filepath.stat().st_size

    info: dict[str, Any] = {
        "format":                  ext,
        "size_bytes":              size_bytes,
        "size_kb":                 round(size_bytes / 1024, 1),
        "page_count":              None,
        "has_extractable_text":    None,
        "is_scanned":              False,
        "has_images":              None,
        "has_tables":              None,
        "sample_text":             "",
        "language":                "en",  # heuristic only; could be improved
    }

    if ext == "pdf":
        try:
            import pdfplumber
            with pdfplumber.open(str(filepath)) as pdf:
                pages = pdf.pages
                info["page_count"] = len(pages)

                # OCR signal across the WHOLE doc (or a sample for huge docs).
                # Critical for mixed-content PDFs where pages 1-3 are
                # born-digital but pages 40+ are scanned — the old "sample
                # first 3 pages" heuristic missed that case entirely.
                ocr_info = _scan_pdf_ocr_signal(pages)
                info["ocr_fraction"]      = ocr_info["ocr_fraction"]
                info["pages_needing_ocr"] = ocr_info["pages_needing_ocr"]
                info["pages_sampled"]     = ocr_info["pages_sampled"]
                info["sampled_indices"]   = ocr_info["sampled_indices"][:30]  # trim for payload size

                # Text sample for the agent's prompt context — first 3 pages
                # are usually enough to recognise the doc's nature.
                sample_parts: list[str] = []
                for p in pages[:3]:
                    try:
                        sample_parts.append((p.extract_text() or "").strip())
                    except Exception:
                        pass
                joined = "\n".join(sample_parts).strip()
                info["sample_text"]          = joined[:1500]
                info["has_extractable_text"] = ocr_info["ocr_fraction"] < 1.0
                # is_scanned now means "enough of this doc needs OCR that
                # pdfplumber alone will leave content on the floor"
                info["is_scanned"]           = ocr_info["ocr_fraction"] > _OCR_FRACTION_THRESH

                # Image / table hints — sample first 5 pages
                try:
                    info["has_images"] = any(p.images for p in pages[:5])
                except Exception:
                    info["has_images"] = False
                try:
                    info["has_tables"] = any(p.find_tables() for p in pages[:5])
                except Exception:
                    info["has_tables"] = False
        except Exception as exc:
            info["inspect_error"] = f"pdfplumber failed: {exc}"
    elif ext in ("docx",):
        try:
            import docx as _docx
            doc = _docx.Document(str(filepath))
            paragraphs = list(doc.paragraphs)
            sample = "\n".join(p.text for p in paragraphs[:30] if p.text.strip())
            info["sample_text"]          = sample[:1500]
            info["has_extractable_text"] = len(sample) > 30
            info["has_tables"]           = len(doc.tables) > 0
            # python-docx doesn't easily report image presence
            info["has_images"]           = None
            info["page_count"]           = None  # DOCX has no fixed pages
        except Exception as exc:
            info["inspect_error"] = f"python-docx failed: {exc}"
    elif ext == "xlsx":
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(filepath), data_only=True, read_only=True)
            info["has_tables"]           = True
            info["has_extractable_text"] = True
            info["sheet_names"]          = wb.sheetnames[:10]
            info["sample_text"]          = f"Workbook with sheets: {', '.join(wb.sheetnames[:5])}"
            wb.close()
        except Exception as exc:
            info["inspect_error"] = f"openpyxl failed: {exc}"
    elif ext == "pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(filepath))
            info["page_count"]           = len(prs.slides)
            info["has_extractable_text"] = True
            info["has_images"]           = True  # PPTX usually has visual content
            sample_parts: list[str] = []
            for slide in list(prs.slides)[:3]:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        sample_parts.append(shape.text)
            info["sample_text"] = "\n".join(sample_parts)[:1500]
        except Exception as exc:
            info["inspect_error"] = f"python-pptx failed: {exc}"
    elif ext in ("html", "htm"):
        try:
            html = filepath.read_text(encoding="utf-8", errors="ignore")
            info["has_extractable_text"] = True
            info["has_tables"]           = "<table" in html.lower()
            info["has_images"]           = "<img" in html.lower()
            from bs4 import BeautifulSoup  # noqa: F401 — pulls in BS4
            soup = BeautifulSoup(html, "html.parser")
            info["sample_text"] = soup.get_text(" ", strip=True)[:1500]
        except Exception as exc:
            info["inspect_error"] = f"html inspect failed: {exc}"

    # Stash for downstream tools that want to read it back
    cache.put(job_id, "_agent_inspect", info)
    return info


async def _parse_pdf_native(filepath: Path, job_id: str, **_) -> dict:
    """Use pdfplumber-based parsing (Mode A stage 3)."""
    from pipelines.custom.stage_02_format_detect import run as detect_run
    from pipelines.custom.stage_03_parser import run as parser_run

    det = await detect_run(filepath)
    mime = (det.payload or {}).get("true_mime", "application/pdf")
    res = await parser_run(filepath, mime)
    payload = res.payload or {}
    text_blocks = payload.get("text_blocks") or []
    tables      = payload.get("tables") or []

    cache.put(job_id, "parser_payload", payload)
    cache.put(job_id, "extracted_tables", tables)
    cache.put(job_id, "_agent_text_blocks", text_blocks)
    return {
        "parser":      payload.get("parser_used", "pdfplumber"),
        "page_count":  payload.get("page_count"),
        "word_count":  payload.get("word_count", 0),
        "table_count": len(tables),
        "text_blocks_count": len(text_blocks),
    }


async def _parse_with_docling(filepath: Path, job_id: str, do_ocr: bool = True, **_) -> dict:
    """Use Docling's unified parse (Mode B stage 2)."""
    from pipelines.docling import stage_02_unified_parse
    res = await stage_02_unified_parse.run(filepath, job_id)
    p = res.payload or {}

    # Docling stage writes d_chunks / d_parser_payload — surface them as the
    # agent's primary outputs by copying to un-prefixed keys (so downstream
    # agent tools read the same keys regardless of which parser ran).
    parser_payload = cache.get(job_id, "d_parser_payload", {})
    cache.put(job_id, "parser_payload",   parser_payload)
    cache.put(job_id, "extracted_tables", parser_payload.get("tables", []))
    cache.put(job_id, "_agent_text_blocks", parser_payload.get("text_blocks") or [])
    # Pre-chunked output also available
    cache.put(job_id, "_agent_chunks", cache.get(job_id, "d_chunks", []))
    cache.put(job_id, "_agent_intel_payload", cache.get(job_id, "d_intel_payload", {}))
    return {
        "parser":            p.get("parser"),
        "page_count":        p.get("page_count"),
        "word_count":        p.get("word_count", 0),
        "table_count":       p.get("table_count", 0),
        "chunk_count":       p.get("chunk_count", 0),
        "ocr_used":          p.get("ocr_used", False),
        "ocr_pages_count":   p.get("ocr_pages_count", 0),
        "doc_type":          p.get("doc_type"),
    }


async def _parse_with_vision_ocr(filepath: Path, job_id: str, **_) -> dict:
    """Vision OCR via Claude — last resort for scanned / CID-broken PDFs.

    Renders each PDF page to a PNG via PyMuPDF then asks Claude vision to
    transcribe text + extract tables. Costs ~$0.01-0.02 per page. Sets up the
    same cache keys as the other parsers so chunk_text / embed_and_index work
    without changes after.
    """
    from pipelines.custom.stage_03_parsers import vision_pdf
    payload = await vision_pdf.vision_ocr_full_pdf(filepath)

    # Mirror the cache shape that other parsers set, so downstream tools read
    # the same keys regardless of which parser ran.
    cache.put(job_id, "parser_payload",   payload)
    cache.put(job_id, "extracted_tables", payload.get("tables", []))
    cache.put(job_id, "_agent_text_blocks", payload.get("text_blocks") or [])

    return {
        "parser":              payload.get("parser_used"),
        "page_count":          payload.get("page_count"),
        "pages_processed":     payload.get("vision_pages_processed"),
        "pages_truncated":     payload.get("vision_pages_truncated"),
        "word_count":          payload.get("word_count", 0),
        "table_count":         len(payload.get("tables", [])),
        "ocr_used":            True,
        "ocr_pages_count":     payload.get("ocr_pages_count", 0),
        "failed_pages":        payload.get("vision_failed_pages") or [],
        "elapsed_s":           payload.get("vision_elapsed_s"),
        "vision_input_tokens": payload.get("vision_input_tokens", 0),
        "vision_output_tokens": payload.get("vision_output_tokens", 0),
    }


async def _parse_office_document(filepath: Path, job_id: str, **_) -> dict:
    """Mode A's parser router handles DOCX/XLSX/PPTX/HTML internally."""
    return await _parse_pdf_native(filepath, job_id)


async def _chunk_text(filepath: Path, job_id: str, **_) -> dict:
    from pipelines.custom import stage_05_chunker

    # Prefer chunks already produced (Docling pre-chunks)
    existing = cache.get(job_id, "_agent_chunks")
    if existing:
        cache.put(job_id, "chunks", existing)
        return {"chunk_count": len(existing), "source": "from_parser"}

    text_blocks = cache.get(job_id, "_agent_text_blocks", [])
    parser_payload = {"text_blocks": text_blocks, "tables": cache.get(job_id, "extracted_tables", [])}
    res = await stage_05_chunker.run(parser_payload)
    chunks = (res.payload or {}).get("chunks", [])
    cache.put(job_id, "chunks", chunks)
    return {
        "chunk_count":           len(chunks),
        "avg_chunk_size_tokens": (res.payload or {}).get("avg_chunk_size_tokens", 0),
        "coverage_pct":          (res.payload or {}).get("coverage_pct", 0),
        "source":                "freshly_chunked",
    }


async def _describe_tables(filepath: Path, job_id: str, **_) -> dict:
    """Tool: describe parser-extracted tables + push table_summary chunks
    into the embedding queue. Idempotent (deduped by chunk id)."""
    from pipelines.custom import stage_06_multimodal
    from pipelines.custom.runner import _table_summary_chunks

    parser_payload = cache.get(job_id, "parser_payload", {}) or {}
    tables = parser_payload.get("tables") or []
    if not tables:
        return {
            "tables_described": 0, "table_summary_chunks_added": 0,
            "reason": "parser found no structured tables in this document",
        }

    result = await stage_06_multimodal.enrich_tables(parser_payload)
    newly_enriched = result.get("tables_enriched", [])

    # APPEND to enriched_tables (don't replace — caption_images may have
    # already added OCR'd tables here). Dedup by `id` so re-runs are idempotent.
    existing_enriched = cache.get(job_id, "enriched_tables", []) or []
    existing_enriched_ids = {t.get("id") for t in existing_enriched}
    to_append = [t for t in newly_enriched if t.get("id") not in existing_enriched_ids]
    combined_enriched = existing_enriched + to_append
    cache.put(job_id, "enriched_tables", combined_enriched)

    # Build summary chunks from the COMBINED list so doc_table_N IDs stay
    # consistent regardless of which tool (describe_tables / caption_images)
    # ran first.
    table_summary_added = 0
    if combined_enriched:
        summary_chunks = _table_summary_chunks(combined_enriched)
        existing = cache.get(job_id, "chunks", []) or []
        existing_ids = {c.get("id") for c in existing}
        new_summary = [c for c in summary_chunks if c.get("id") not in existing_ids]
        if new_summary:
            cache.put(job_id, "chunks", existing + new_summary)
            table_summary_added = len(new_summary)

    return {
        "tables_described":           result.get("tables_described", 0),
        "table_summary_chunks_added": table_summary_added,
        "llm_input_tokens":           result.get("llm_input_tokens", 0),
        "llm_output_tokens":          result.get("llm_output_tokens", 0),
        "llm_cost_usd":               result.get("llm_cost_usd", 0.0),
    }


async def _caption_images(filepath: Path, job_id: str, **_) -> dict:
    """Tool: image captioning + structured OCR of scanned pages.

    Routes the OCR'd content correctly:
    - OCR'd prose chunks → chunks cache (for vector embedding)
    - OCR'd tables → extracted_tables (for SQL store)
                   → enriched_tables (already have descriptions from OCR call)
                   → chunks (table_summary chunks for vector retrieval)

    OCR'd tables come with descriptions already (the structured OCR prompt
    extracts both text AND a one-sentence description per table), so no
    additional Claude calls are needed to summarize them.
    """
    from pipelines.custom import stage_06_multimodal
    from pipelines.custom.runner import _table_summary_chunks

    parser_payload = cache.get(job_id, "parser_payload", {}) or {}
    images = [img for img in (parser_payload.get("images") or []) if img.get("bytes_b64")]
    if not images:
        return {
            "images_captioned": 0, "ocr_pages": 0,
            "reason": "parser found no images in this document",
        }

    result = await stage_06_multimodal.process_images(parser_payload)

    # OCR'd prose chunks → chunks cache (for vector embedding)
    if result.get("ocr_chunks"):
        existing = cache.get(job_id, "chunks", []) or []
        cache.put(job_id, "chunks", existing + result["ocr_chunks"])

    # OCR'd tables → THREE caches in lockstep:
    #   1. extracted_tables — for SQL store (`store_tables_sql` reads this)
    #   2. enriched_tables  — already have descriptions from OCR; no extra Claude
    #   3. chunks           — table_summary chunks for vector retrieval
    ocr_table_summary_added = 0
    if result.get("ocr_tables"):
        ocr_tables = result["ocr_tables"]

        # (1) SQL — append to extracted_tables
        existing_extracted = cache.get(job_id, "extracted_tables", []) or []
        existing_extracted_ids = {t.get("id") for t in existing_extracted}
        new_extracted = [t for t in ocr_tables if t.get("id") not in existing_extracted_ids]
        if new_extracted:
            cache.put(job_id, "extracted_tables", existing_extracted + new_extracted)

        # (2) Enriched — append (OCR tables already carry descriptions)
        existing_enriched = cache.get(job_id, "enriched_tables", []) or []
        existing_enriched_ids = {t.get("id") for t in existing_enriched}
        new_enriched = [t for t in ocr_tables if t.get("id") not in existing_enriched_ids]
        combined_enriched = existing_enriched + new_enriched
        if new_enriched:
            cache.put(job_id, "enriched_tables", combined_enriched)

        # (3) Vector DB — build summary chunks from the COMBINED enriched list
        # (so doc_table_N IDs match what describe_tables would produce)
        if combined_enriched:
            summary_chunks = _table_summary_chunks(combined_enriched)
            existing_chunks = cache.get(job_id, "chunks", []) or []
            existing_chunk_ids = {c.get("id") for c in existing_chunks}
            new_summaries = [c for c in summary_chunks if c.get("id") not in existing_chunk_ids]
            if new_summaries:
                cache.put(job_id, "chunks", existing_chunks + new_summaries)
                ocr_table_summary_added = len(new_summaries)

    return {
        "images_captioned":              len(result.get("captions", [])),
        "captions_failed":               result.get("captions_failed", 0),
        "ocr_pages":                     result.get("ocr_pages_count", 0),
        "ocr_pages_failed":              result.get("ocr_pages_failed", 0),
        "ocr_tables_extracted":          len(result.get("ocr_tables", [])),
        "ocr_table_summary_chunks_added": ocr_table_summary_added,
        "llm_input_tokens":              result.get("llm_input_tokens", 0),
        "llm_output_tokens":             result.get("llm_output_tokens", 0),
        "llm_cost_usd":                  result.get("llm_cost_usd", 0.0),
    }


async def _embed_and_index(filepath: Path, job_id: str, **_) -> dict:
    from pipelines.custom import stage_07_embedding, stage_09_vector_store

    chunks = cache.get(job_id, "chunks", []) or []
    if not chunks:
        return {"error": "no chunks to embed — call chunk_text first"}

    # Embedding
    emb_res = await stage_07_embedding.run(chunks, job_id)
    emb_p = emb_res.payload or {}

    # Vector store (Qdrant + chunk_breakdown side-effect)
    vs_res = await stage_09_vector_store.run(
        job_id, collection_prefix="rag_proto_agent", cache_prefix="",
    )
    vs_p = vs_res.payload or {}

    return {
        "chunks_embedded":     emb_p.get("chunks_embedded", 0),
        "vector_dim":          emb_p.get("vector_dim"),
        "model":               emb_p.get("model"),
        "use_real_embeddings": emb_p.get("use_real_embeddings", False),
        "embedding_cost_usd":  emb_p.get("llm_cost_usd", 0.0),
        "vectors_upserted":    vs_p.get("vectors_upserted", 0),
        "qdrant_live":         vs_p.get("qdrant_live", False),
        "collection":          vs_p.get("collection"),
    }


async def _store_tables_sql(filepath: Path, job_id: str, **_) -> dict:
    from services import sql_store
    extracted = cache.get(job_id, "extracted_tables", []) or []
    if not extracted:
        return {"sql_tables_created": 0, "reason": "no extracted_tables in cache"}
    registry = sql_store.create_tables(extracted, job_id, cache_prefix="")
    cache.put(job_id, "sql_registry", registry)
    return {
        "sql_tables_created": len(registry),
        "tables": [
            {"name": k, "row_count": v["row_count"], "columns": v["original_headers"][:5]}
            for k, v in registry.items()
        ],
    }


async def _extract_entities(filepath: Path, job_id: str, **_) -> dict:
    from pipelines.custom import stage_09_knowledge_graph
    res = await stage_09_knowledge_graph.run(job_id, cache_prefix="")
    p = res.payload or {}
    return {
        "entity_count":       p.get("entity_count", 0),
        "relationship_count": p.get("relationship_count", 0),
        "entity_types":       p.get("unique_entity_types", []),
        "top_entities":       (p.get("top_entities") or [])[:5],
    }


async def _finalize(filepath: Path, job_id: str, summary: str = "", **_) -> dict:
    return {"final": True, "summary": summary}


# ── Dispatch table ────────────────────────────────────────────────────────────

_TOOL_EXECUTORS: dict[str, Callable[..., Awaitable[dict]]] = {
    "inspect_document":     _inspect_document,
    "parse_pdf_native":     _parse_pdf_native,
    "parse_with_docling":   _parse_with_docling,
    "parse_with_vision_ocr": _parse_with_vision_ocr,
    "parse_office_document": _parse_office_document,
    "chunk_text":           _chunk_text,
    "describe_tables":      _describe_tables,
    "caption_images":       _caption_images,
    "embed_and_index":      _embed_and_index,
    "store_tables_sql":     _store_tables_sql,
    "extract_entities":     _extract_entities,
    "finalize":             _finalize,
}


async def execute_tool(
    name: str,
    tool_input: dict,
    filepath: Path,
    job_id: str,
) -> dict:
    """Execute a tool by name with the agent's input + the implicit job context.

    Tools always have `filepath` and `job_id` available. Any extra agent-
    supplied args are passed as kwargs.
    """
    fn = _TOOL_EXECUTORS.get(name)
    if fn is None:
        return {"error": f"unknown tool: {name}"}
    t0 = time.perf_counter()
    try:
        result = await fn(filepath=filepath, job_id=job_id, **(tool_input or {}))
        result["_elapsed_ms"] = round((time.perf_counter() - t0) * 1000, 1)
        return result
    except Exception as exc:
        log.exception("tool %s failed", name)
        return {
            "error":       f"{type(exc).__name__}: {exc}",
            "_elapsed_ms": round((time.perf_counter() - t0) * 1000, 1),
        }
