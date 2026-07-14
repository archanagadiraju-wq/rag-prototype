"""Generate realistic demo documents in demo_docs/.

Run once:  python create_demo_docs.py
"""
from __future__ import annotations
import io
from pathlib import Path

import openpyxl
from openpyxl.styles import Font as XlFont, PatternFill, Alignment
import docx
from docx.shared import Pt
import pptx
from pptx.util import Inches, Pt as PptPt
from pptx.dml.color import RGBColor

DEMO_DOCS_DIR = Path(__file__).parent.parent / "demo_docs"
DEMO_DOCS_DIR.mkdir(parents=True, exist_ok=True)


# ── Multi-page text-native PDF ─────────────────────────────────────────────────

def _make_multipage_pdf(pages: list[list[str]]) -> bytes:
    """Build a multi-page text-native PDF with correct xref offsets."""
    out = bytearray(b"%PDF-1.4\n")
    offsets: dict[int, int] = {}

    n = len(pages)
    # obj numbers: 1=catalog, 2=pages, 3..2+n=page objs, 3+n..2+2n=streams, 3+2n=font
    catalog_n = 1
    pages_n   = 2
    page_ns   = list(range(3, 3 + n))
    stream_ns = list(range(3 + n, 3 + 2 * n))
    font_n    = 3 + 2 * n

    def w(obj_num: int, body: bytes) -> None:
        offsets[obj_num] = len(out)
        out.extend(f"{obj_num} 0 obj\n".encode())
        out.extend(body)
        out.extend(b"\nendobj\n\n")

    kids = " ".join(f"{p} 0 R" for p in page_ns)
    w(catalog_n, f"<< /Type /Catalog /Pages {pages_n} 0 R >>".encode())
    w(pages_n,   f"<< /Type /Pages /Kids [{kids}] /Count {n} >>".encode())

    for pn, sn in zip(page_ns, stream_ns):
        w(pn, (
            f"<< /Type /Page /Parent {pages_n} 0 R /MediaBox [0 0 612 792]\n"
            f"   /Contents {sn} 0 R /Resources << /Font << /F1 {font_n} 0 R "
            f"/F2 {font_n + 1} 0 R >> >> >>"
        ).encode())

    for sn, lines in zip(stream_ns, pages):
        parts = ["BT", "/F1 11 Tf", "50 740 Td", "14 TL"]
        for line in lines:
            safe = line.replace("\\","\\\\").replace("(","\\(").replace(")","\\)")
            if line.startswith("##"):
                parts += ["/F2 13 Tf", f"({safe[2:].strip()}) Tj", "T*", "/F1 11 Tf"]
            else:
                parts.append(f"({safe}) Tj")
                parts.append("T*")
        parts.append("ET")
        s = "\n".join(parts).encode()
        w(sn, f"<< /Length {len(s)} >>\nstream\n".encode() + s + b"\nendstream")

    w(font_n,     b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")
    w(font_n + 1, b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica-Bold >>")

    total = font_n + 1
    xref_off = len(out)
    out.extend(f"xref\n0 {total + 1}\n0000000000 65535 f \n".encode())
    for i in range(1, total + 1):
        out.extend(f"{offsets[i]:010d} 00000 n \n".encode())
    out.extend((
        f"trailer\n<< /Size {total + 1} /Root {catalog_n} 0 R >>\n"
        f"startxref\n{xref_off}\n%%EOF\n"
    ).encode())
    return bytes(out)


# ── 01 Pharmaceutical trial PDF ────────────────────────────────────────────────

def make_pharma_pdf() -> bytes:
    p1 = [
        "## XR-2041 Phase III Randomised Controlled Trial",
        "A Multicentre, Double-Blind, Placebo-Controlled Study of XR-2041",
        "in Patients with Moderate-to-Severe Rheumatoid Arthritis",
        "",
        "Authors: Chen W, Patel R, Okonkwo A, Vasquez M, Nakamura T",
        "Sponsor: Helix Therapeutics Inc.  |  Protocol: HX-2041-P3-001",
        "Submission Date: 14 March 2026",
        "",
        "## Abstract",
        "Background: Rheumatoid arthritis (RA) affects approximately 1% of the global",
        "population and is characterised by synovial inflammation and joint destruction.",
        "Current biologics achieve remission in only 30-40% of patients, leaving a",
        "significant unmet need for novel therapies.",
        "",
        "Objective: To evaluate the efficacy and safety of XR-2041, a selective JAK1/JAK3",
        "inhibitor, compared with placebo in adults with inadequate response to MTX.",
        "",
        "Design: 52-week, Phase III, randomised (2:1), double-blind, placebo-controlled trial",
        "conducted across 48 sites in 12 countries.",
        "",
        "Participants: 1,248 adults (18-75 yr) with ACR/EULAR-defined moderate-to-severe RA,",
        "DAS28-CRP >= 3.2, and documented inadequate response to methotrexate (MTX).",
        "",
        "Interventions: XR-2041 15 mg once daily (n=832) or matched placebo (n=416),",
        "both with stable background MTX 7.5-25 mg/week.",
        "",
        "Primary Endpoint: ACR20 response rate at Week 12.",
        "Key Secondary Endpoints: DAS28-CRP remission, ACR50/70 at Week 12 and 52,",
        "HAQ-DI change from baseline, radiographic progression (mTSS) at Week 52.",
        "",
        "Results: ACR20 at Week 12: XR-2041 70.4% vs placebo 29.1% (p<0.001, OR 5.82).",
        "DAS28-CRP remission at Week 12: 38.2% vs 8.4% (p<0.001).",
        "ACR50 at Week 12: 48.7% vs 14.1%. ACR70 at Week 12: 26.3% vs 5.3%.",
        "No new safety signals; serious adverse events: 6.1% XR-2041 vs 5.8% placebo.",
        "",
        "Conclusion: XR-2041 demonstrated superior efficacy vs placebo across all primary",
        "and key secondary endpoints with an acceptable safety profile.",
    ]
    p2 = [
        "## 1. Introduction",
        "Rheumatoid arthritis (RA) is a systemic autoimmune disease primarily affecting",
        "synovial joints, with a worldwide prevalence of 0.5-1.0%. The pathogenesis",
        "involves dysregulated cytokine signalling — particularly IL-6, TNF-alpha, and",
        "IL-17 — mediated through Janus kinase (JAK) pathways.",
        "",
        "Despite the availability of conventional DMARDs and biologic agents (TNFi, IL-6Ri,",
        "anti-CD20), approximately 30-40% of patients fail to achieve sustained remission.",
        "JAK inhibitors represent the latest class of targeted synthetic DMARDs (tsDMARDs),",
        "with tofacitinib, baricitinib, and upadacitinib approved for RA.",
        "",
        "XR-2041 is a next-generation oral JAK1/JAK3 inhibitor with high selectivity",
        "(JAK1 IC50 = 3.2 nM; JAK3 IC50 = 5.8 nM; >200-fold selectivity over JAK2),",
        "designed to minimise off-target haematological effects associated with JAK2",
        "inhibition observed with earlier agents.",
        "",
        "## 2. Methods",
        "## 2.1 Study Design",
        "This was a global, multicentre, randomised, double-blind, placebo-controlled,",
        "parallel-group Phase III study. Patients were randomised 2:1 to XR-2041 or",
        "placebo using interactive response technology (IRT), stratified by:",
        "  - Geographic region (North America / Europe / Asia-Pacific / Other)",
        "  - Prior biologic use (naive vs experienced)",
        "  - Baseline DAS28-CRP (<5.1 vs >=5.1)",
        "",
        "## 2.2 Participants",
        "Inclusion: Age 18-75; ACR/EULAR 2010 RA classification criteria; DAS28-CRP >=3.2;",
        "stable MTX (>=12 weeks, 7.5-25 mg/week); inadequate response defined as",
        ">=3 months of therapy without remission.",
        "",
        "Exclusion: Prior JAK inhibitor use; active or latent TB (QuantiFERON positive);",
        "eGFR <30 mL/min/1.73m2; ALT/AST >2x ULN; current malignancy; pregnancy.",
        "",
        "## 2.3 Interventions",
        "XR-2041 15 mg film-coated tablet, once daily, orally.",
        "Placebo: matched tablet, once daily, orally.",
        "Background: stable MTX throughout; folic acid 5 mg/week mandatory.",
        "Rescue: IV methylprednisolone 125 mg single dose permitted once for flare.",
    ]
    p3 = [
        "## 3. Results",
        "## 3.1 Patient Disposition and Baseline",
        "Screened: 1,847  |  Randomised: 1,248  |  Completed 52 weeks: 1,089 (87.3%)",
        "XR-2041: n=832 (completed=731, 87.9%)  |  Placebo: n=416 (completed=358, 86.1%)",
        "",
        "Baseline characteristics were well balanced:",
        "  Mean age: 51.3 yr (XR) vs 51.8 yr (PBO)",
        "  Female: 78.1% vs 77.6%",
        "  Mean DAS28-CRP: 5.81 vs 5.79",
        "  Mean HAQ-DI: 1.48 vs 1.51",
        "  Biologic-naive: 71.4% vs 70.9%",
        "  Median RA duration: 7.2 yr vs 7.5 yr",
        "  Mean CRP: 18.4 mg/L vs 19.1 mg/L",
        "  Mean SJC28: 12.3 vs 12.6",
        "",
        "## 3.2 Primary Endpoint",
        "ACR20 Response at Week 12:",
        "  XR-2041: 586/832 (70.4%)  vs  Placebo: 121/416 (29.1%)",
        "  Odds Ratio: 5.82 (95% CI 4.31-7.86)  |  p-value: <0.0001",
        "  Risk Difference: 41.3% (95% CI 35.9-46.7%)",
        "  NNT: 2.4 (95% CI 2.1-2.8)",
        "",
        "## 3.3 Key Secondary Endpoints at Week 12",
        "  DAS28-CRP <2.6 (remission): 38.2% vs 8.4%   (OR 6.65; p<0.0001)",
        "  DAS28-CRP <=3.2 (LDA):      56.7% vs 17.8%  (OR 5.97; p<0.0001)",
        "  ACR50 response:              48.7% vs 14.1%  (OR 5.67; p<0.0001)",
        "  ACR70 response:              26.3% vs  5.3%  (OR 6.31; p<0.0001)",
        "  CDAI remission (<=2.8):      22.1% vs  4.1%  (OR 6.60; p<0.0001)",
        "  HAQ-DI change from BL:      -0.61 vs -0.22  (diff -0.39; p<0.0001)",
        "  Patient Global VAS change:  -31.4 vs -12.8  (diff -18.6; p<0.0001)",
        "",
        "## 3.4 Radiographic Outcomes at Week 52",
        "  mTSS change: 0.42 (XR) vs 1.87 (PBO) [diff -1.45; 95%CI -2.01/-0.89; p<0.0001]",
        "  Erosion score change: 0.18 vs 0.94 (diff -0.76; p<0.0001)",
        "  JSN score change: 0.24 vs 0.93 (diff -0.69; p<0.0001)",
        "  No radiographic progression (deltaTSS<=0): 68.4% vs 44.7% (p<0.0001)",
    ]
    p4 = [
        "## 3.5 Safety Summary",
        "Treatment-emergent AEs (TEAEs): 74.3% XR-2041 vs 66.8% placebo",
        "Serious TEAEs: 6.1% vs 5.8%  |  TEAE leading to discontinuation: 3.2% vs 2.4%",
        "",
        "Most common TEAEs (>=5% in either arm):",
        "  Upper respiratory tract infection: 14.8% vs 10.6%",
        "  Nasopharyngitis:  9.2% vs  7.9%",
        "  Headache:         8.1% vs  6.7%",
        "  Nausea:           6.4% vs  3.8%",
        "  Diarrhoea:        5.3% vs  4.1%",
        "  Urinary tract infection: 5.1% vs  3.6%",
        "",
        "Events of special interest:",
        "  Serious infections: 2.4% vs 2.2%  (incidence rate 2.8 vs 2.5/100 PY)",
        "  Herpes zoster: 1.8% vs 0.7%  (IR 2.1 vs 0.8/100 PY)",
        "  MACE: 0.2% vs 0.0%  (IR 0.24/100 PY; adjudicated by independent committee)",
        "  DVT/PE: 0.4% vs 0.2%  (all patients had identifiable risk factors)",
        "  Malignancies: 0.6% vs 0.5%  (consistent with RA background rates)",
        "  Laboratory: mean Hgb -0.3 g/dL; ANC stable; no Grade 4 cytopenias",
        "  Creatinine: mean +0.04 mg/dL from baseline, stable throughout",
        "",
        "## 4. Discussion",
        "XR-2041 15 mg once daily met its primary endpoint with an ACR20 response of 70.4%",
        "at Week 12, versus 29.1% for placebo (p<0.0001). This effect size is consistent",
        "with or exceeds approved JAK inhibitors in similar RA populations.",
        "",
        "The selectivity profile of XR-2041 (>200-fold JAK1/3 over JAK2) was associated",
        "with a stable haematological profile: no Grade 3/4 anaemia, neutropenia, or",
        "thrombocytopenia events were observed. MACE and VTE rates were low and consistent",
        "with the expected background rate for an RA population of this age and CV profile.",
        "",
        "Limitations include the 12-week primary endpoint (regulatory minimum) and the",
        "exclusion of patients with eGFR <30, limiting generalisability to severe CKD.",
        "",
        "## 5. Conclusion",
        "XR-2041 15 mg once daily demonstrated statistically significant and clinically",
        "meaningful superiority over placebo on ACR20 at Week 12 and across all key",
        "secondary endpoints including DAS28-CRP remission, ACR50/70, HAQ-DI improvement,",
        "and radiographic non-progression at Week 52. The safety profile was acceptable",
        "with no new signals beyond the class effect of JAK inhibitors.",
        "These results support XR-2041 as a potential new treatment option for patients",
        "with moderate-to-severe RA with inadequate response to methotrexate.",
    ]
    return _make_multipage_pdf([p1, p2, p3, p4])


# ── 02 Financial model XLSX ────────────────────────────────────────────────────

def make_financial_xlsx() -> bytes:
    wb = openpyxl.Workbook()

    # ── Revenue sheet ──
    ws = wb.active
    ws.title = "Revenue"
    hdr = ["Month", "New MRR ($)", "Expansion MRR ($)", "Churned MRR ($)", "Net New MRR ($)",
           "Total MRR ($)", "ARR ($)", "Customers", "ARPU ($)", "Growth MoM (%)"]
    ws.append(hdr)
    mrr, customers = 180_000, 42
    months = ["Jan-25","Feb-25","Mar-25","Apr-25","May-25","Jun-25",
              "Jul-25","Aug-25","Sep-25","Oct-25","Nov-25","Dec-25",
              "Jan-26","Feb-26","Mar-26","Apr-26","May-26","Jun-26",
              "Jul-26","Aug-26","Sep-26","Oct-26","Nov-26","Dec-26"]
    for m in months:
        new = round(mrr * 0.18)
        exp = round(mrr * 0.04)
        churn = round(mrr * 0.015)
        net_new = new + exp - churn
        mrr += net_new
        customers += 3
        arpu = round(mrr / customers)
        growth = round(net_new / (mrr - net_new) * 100, 1)
        ws.append([m, new, exp, churn, net_new, mrr, mrr*12, customers, arpu, growth])

    # ── Assumptions sheet ──
    ws2 = wb.create_sheet("Assumptions")
    ws2.append(["Driver", "Value", "Notes"])
    for row in [
        ("Monthly Growth Rate (target)", "12%", "Board-approved plan"),
        ("Gross Margin (%)", "74%", "Incl. hosting, support, COGS"),
        ("Net Revenue Retention (%)", "118%", "Based on last 4Q average"),
        ("Sales Efficiency Ratio", "0.82", "ARR added / S&M spend"),
        ("Payback Period (months)", "14", "CAC payback on gross margin basis"),
        ("Avg Contract Length (months)", "18", "Mix of monthly/annual"),
        ("Avg Deal Size (ACV)", "52,000", "New logo average"),
        ("Churn Rate (monthly)", "1.5%", "Gross revenue churn"),
        ("Expansion Rate (monthly)", "4.0%", "Upsell + cross-sell"),
        ("R&D as % of Revenue", "22%", "Headcount + tooling"),
        ("S&M as % of Revenue", "28%", "Full-cycle including SDR"),
        ("G&A as % of Revenue", "9%", "Finance, legal, HR"),
    ]:
        ws2.append(row)

    # ── Headcount sheet ──
    ws3 = wb.create_sheet("Headcount")
    ws3.append(["Department", "Q1-25", "Q2-25", "Q3-25", "Q4-25",
                "Q1-26", "Q2-26", "Q3-26", "Q4-26", "Annual Cost ($M)"])
    for dept, q1,q2,q3,q4,q1b,q2b,q3b,q4b,cost in [
        ("Engineering",   18,20,24,27, 30,33,36,40, 6.8),
        ("Product",        4, 5, 6, 7,  8, 9,10,11, 1.9),
        ("Sales",          8,10,13,16, 18,21,24,27, 4.2),
        ("Marketing",      4, 5, 6, 7,  8, 9,10,11, 1.6),
        ("Customer Success",5, 6, 7, 8, 10,11,13,14, 1.8),
        ("G&A",            4, 4, 5, 5,  6, 6, 7, 8, 1.2),
        ("Total",         43,50,61,70, 80,89,100,111,17.5),
    ]:
        ws3.append([dept,q1,q2,q3,q4,q1b,q2b,q3b,q4b,cost])

    # ── P&L sheet ──
    ws4 = wb.create_sheet("P&L")
    ws4.append(["Line Item", "Q1-25","Q2-25","Q3-25","Q4-25","FY-25","Q1-26","Q2-26","FY-26E"])
    for row in [
        ("Revenue",        620,730,860,1020,3230,1190,1390,5800),
        ("COGS",          -161,-190,-224,-265,-840,-310,-362,-1508),
        ("Gross Profit",   459,540,636, 755,2390, 880,1028,4292),
        ("Gross Margin %", "74%","74%","74%","74%","74%","74%","74%","74%"),
        ("R&D",           -136,-161,-189,-224,-710,-262,-306,-1276),
        ("S&M",           -174,-204,-241,-286,-905,-333,-389,-1624),
        ("G&A",            -56, -66, -77, -92,-291,-107,-125, -522),
        ("Op. Loss",        93, 109, 129,  153, 484,  178,  208,  870),
        ("EBITDA Margin%","15%","15%","15%","15%","15%","15%","15%","15%"),
    ]:
        ws4.append(list(row))

    # ── Scenarios sheet ──
    ws5 = wb.create_sheet("Scenarios")
    ws5.append(["Metric", "Bear (-30%)", "Base", "Bull (+25%)", "Bull Driver"])
    for row in [
        ("ARR EOY 2026 ($M)",    52, 74, 93, "Enterprise expansion"),
        ("Monthly Growth Rate",  "8%","12%","15%", "New logo acceleration"),
        ("NRR (%)",             "108%","118%","128%", "Product-led expansion"),
        ("Gross Margin (%)",    "70%","74%","77%", "Infra efficiencies"),
        ("Cash Runway (months)", 14, 22, 31, "Burn reduction"),
        ("Headcount EOY",        78, 111, 138, "GTM hiring"),
        ("Series B Target ($M)", 35, 50, 65, "Valuation premium"),
    ]:
        ws5.append(list(row))

    ws6 = wb.create_sheet("Waterfall")
    ws6.append(["Component", "ARR Impact ($M)", "Cumulative ($M)"])
    for row in [
        ("Opening ARR (Jan 2025)",   18.4, 18.4),
        ("New Logo",                 24.8, 43.2),
        ("Expansion (upsell/cross)", 11.2, 54.4),
        ("Gross Churn",              -5.1, 49.3),
        ("Contraction",              -1.4, 47.9),
        ("FX Impact",                -0.8, 47.1),
        ("Closing ARR (Dec 2025)",   47.1, 47.1),
    ]:
        ws6.append(list(row))

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ── 03 Vendor contract DOCX ────────────────────────────────────────────────────

def make_contract_docx() -> bytes:
    doc = docx.Document()
    doc.add_heading("ENTERPRISE SOFTWARE SUBSCRIPTION AGREEMENT", 0)
    doc.add_paragraph(
        "This Enterprise Software Subscription Agreement (\"Agreement\") is entered into as of "
        "the Effective Date set forth below, by and between Acme Corporation, a Delaware "
        "corporation with its principal place of business at 100 Innovation Drive, San Francisco, "
        "CA 94105 (\"Customer\"), and Nexus Software Inc., a Delaware corporation with its "
        "principal place of business at 200 Platform Way, Austin, TX 78701 (\"Vendor\")."
    )
    sections = [
        ("1. DEFINITIONS", [
            ("1.1 \"Authorized Users\"", "means the Customer's employees, contractors, and agents who are authorised by Customer to access and use the Software under this Agreement, up to the number of seats specified in the applicable Order Form."),
            ("1.2 \"Documentation\"", "means Vendor's then-current user manuals, technical specifications, and online help files for the Software, as updated by Vendor from time to time."),
            ("1.3 \"Order Form\"", "means the order document executed by both parties that specifies the Software, subscription tier, number of Authorized Users, Fees, and subscription term, and that is incorporated herein by reference."),
            ("1.4 \"Professional Services\"", "means the implementation, configuration, training, and consulting services described in a Statement of Work (\"SOW\") mutually executed by the parties."),
            ("1.5 \"Software\"", "means Vendor's proprietary cloud-based software platform known as NexusIQ, including all updates, upgrades, and new releases provided by Vendor during the Subscription Term."),
            ("1.6 \"Subscription Term\"", "means the initial period specified in the Order Form and any renewal periods thereafter."),
        ]),
        ("2. SUBSCRIPTION AND ACCESS", [
            ("2.1 Grant of Subscription", "Subject to the terms and conditions of this Agreement and timely payment of Fees, Vendor grants Customer a non-exclusive, non-transferable, non-sublicensable subscription right during the Subscription Term for Authorized Users to access and use the Software solely for Customer's internal business operations."),
            ("2.2 Restrictions", "Customer shall not: (a) sublicense, resell, or distribute the Software; (b) modify, adapt, or create derivative works; (c) reverse engineer, decompile, or disassemble the Software; (d) access the Software to build a competitive product; (e) exceed the number of Authorised Users specified in the applicable Order Form."),
            ("2.3 Vendor Responsibilities", "Vendor shall: (a) maintain the security and integrity of Customer data; (b) provide the Software in conformance with the Documentation; (c) implement commercially reasonable security measures including SOC 2 Type II controls."),
        ]),
        ("3. FEES AND PAYMENT", [
            ("3.1 Subscription Fees", "Customer shall pay the Fees specified in the Order Form. Annual subscriptions are invoiced on the Effective Date and each anniversary thereof. Monthly subscriptions are invoiced at the start of each calendar month."),
            ("3.2 Payment Terms", "All invoices are due and payable within thirty (30) days of invoice date (Net-30). Overdue amounts accrue interest at 1.5% per month or the maximum permitted by law, whichever is lower."),
            ("3.3 Taxes", "Fees exclude all applicable taxes. Customer is responsible for all sales, use, VAT, GST, and similar taxes, excluding taxes on Vendor's income."),
            ("3.4 Price Adjustments", "Vendor may adjust Fees at renewal upon sixty (60) days' written notice, with increases capped at the greater of (a) 5% or (b) the year-over-year change in the US CPI-U for the preceding 12 months."),
        ]),
        ("4. SERVICE LEVEL AGREEMENT", [
            ("4.1 Uptime Commitment", "Vendor commits to 99.9% monthly uptime for the Software (excluding Scheduled Maintenance and events beyond Vendor's reasonable control). \"Uptime\" is calculated as: ((Total Minutes - Downtime Minutes) / Total Minutes) x 100."),
            ("4.2 Service Credits", "If monthly uptime falls below 99.9%, Customer is eligible for service credits: 99.0-99.9% = 10% of monthly fee; 95.0-98.9% = 25%; below 95% = 50%. Credits are Customer's sole remedy for uptime failures."),
            ("4.3 Support", "Vendor provides Tier-1 through Tier-3 technical support via email, chat, and phone (for Enterprise tier). Initial response SLAs: Critical (P1) 1 hour; High (P2) 4 hours; Medium (P3) 1 business day; Low (P4) 3 business days."),
            ("4.4 Scheduled Maintenance", "Vendor will provide 72 hours' advance notice for scheduled maintenance windows. Emergency maintenance may be performed with shorter notice. Maintenance windows shall not exceed 4 hours per month."),
        ]),
        ("5. CONFIDENTIALITY", [
            ("5.1 Obligations", "Each party (\"Receiving Party\") agrees to: (a) hold in strict confidence the Confidential Information of the other party (\"Disclosing Party\"); (b) not disclose such information to third parties without prior written consent; (c) use Confidential Information solely for purposes of this Agreement; (d) restrict access to employees and contractors with a need-to-know who are bound by obligations no less restrictive than those herein."),
            ("5.2 Exclusions", "Confidentiality obligations do not apply to information that: (a) is or becomes publicly known without breach of this Agreement; (b) was rightfully known before disclosure; (c) is independently developed without use of Confidential Information; (d) is required to be disclosed by law, provided Receiving Party gives prompt written notice."),
            ("5.3 Term", "Confidentiality obligations survive for three (3) years following expiration or termination of this Agreement, provided that trade secrets shall be maintained in confidence indefinitely."),
        ]),
        ("6. DATA PROTECTION AND SECURITY", [
            ("6.1 Data Processing", "To the extent Vendor processes personal data on behalf of Customer, Vendor acts as a data processor under applicable privacy laws (including GDPR and CCPA). The parties shall execute Vendor's Data Processing Addendum (\"DPA\"), which is incorporated herein."),
            ("6.2 Security Measures", "Vendor maintains an information security programme that includes: AES-256 encryption at rest; TLS 1.2+ in transit; SOC 2 Type II certification; annual penetration testing by an accredited third party; 24/7 security monitoring."),
            ("6.3 Data Residency", "Customer data is stored in the AWS us-east-1 and eu-west-1 regions by default. Alternative regions available upon written request and subject to additional fees."),
            ("6.4 Breach Notification", "Vendor shall notify Customer within 72 hours of confirming a security incident affecting Customer data, as required by applicable law."),
        ]),
        ("7. INTELLECTUAL PROPERTY", [
            ("7.1 Vendor IP", "Vendor retains all right, title, and interest in the Software, Documentation, and any improvements or modifications thereto. Nothing in this Agreement transfers ownership of Vendor IP to Customer."),
            ("7.2 Customer Data", "Customer retains all right, title, and interest in Customer data. Customer grants Vendor a limited licence to process Customer data solely to provide the Software and related services."),
            ("7.3 Feedback", "If Customer provides Vendor with feedback or suggestions, Vendor may use such feedback without restriction or obligation to Customer."),
        ]),
        ("8. TERM AND TERMINATION", [
            ("8.1 Initial Term", "This Agreement commences on the Effective Date and continues for the Initial Term specified in the Order Form, unless earlier terminated in accordance with this Section."),
            ("8.2 Renewal", "The Agreement automatically renews for successive one-year periods unless either party provides written notice of non-renewal at least sixty (60) days prior to the end of the then-current term."),
            ("8.3 Termination for Cause", "Either party may terminate this Agreement upon thirty (30) days' written notice if the other party materially breaches any term of this Agreement and fails to cure such breach within such period."),
            ("8.4 Effect of Termination", "Upon termination, Customer's right to access the Software ceases. Vendor will make Customer data available for export for thirty (30) days, after which Vendor may delete it. Fees paid are non-refundable except for prepaid unused service following termination for cause."),
        ]),
        ("9. LIMITATION OF LIABILITY", [
            ("9.1 Exclusion of Consequential Damages", "NEITHER PARTY SHALL BE LIABLE FOR INDIRECT, INCIDENTAL, SPECIAL, CONSEQUENTIAL, PUNITIVE, OR EXEMPLARY DAMAGES, INCLUDING LOSS OF PROFITS, REVENUE, DATA, OR BUSINESS OPPORTUNITIES, EVEN IF ADVISED OF THE POSSIBILITY OF SUCH DAMAGES."),
            ("9.2 Cap on Liability", "EACH PARTY'S TOTAL CUMULATIVE LIABILITY ARISING OUT OF OR RELATED TO THIS AGREEMENT SHALL NOT EXCEED THE TOTAL FEES PAID OR PAYABLE BY CUSTOMER IN THE TWELVE (12) MONTHS PRECEDING THE CLAIM."),
            ("9.3 Exceptions", "The foregoing limitations do not apply to: (a) breaches of confidentiality obligations; (b) indemnification obligations; (c) gross negligence or wilful misconduct; (d) death or personal injury caused by negligence."),
        ]),
    ]
    for heading, subsections in sections:
        doc.add_heading(heading, level=1)
        for sub_heading, text in subsections:
            doc.add_heading(sub_heading, level=2)
            doc.add_paragraph(text)

    doc.add_heading("10. GENERAL PROVISIONS", level=1)
    doc.add_paragraph(
        "10.1 Governing Law. This Agreement is governed by the laws of the State of Delaware "
        "without regard to its conflict of law provisions. 10.2 Dispute Resolution. The parties "
        "agree to attempt to resolve disputes through good-faith negotiation for 30 days before "
        "resorting to binding arbitration under JAMS rules in San Francisco, CA. 10.3 Force "
        "Majeure. Neither party is liable for delays caused by circumstances beyond its reasonable "
        "control. 10.4 Entire Agreement. This Agreement, including all Order Forms and SOWs, "
        "constitutes the entire agreement between the parties and supersedes all prior agreements."
    )
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── 04 Technical spec HTML ─────────────────────────────────────────────────────

def make_technical_html() -> bytes:
    return """<!DOCTYPE html>
<html lang="en">
<head><meta charset="UTF-8"><title>NexusIQ REST API v3.2 — Technical Reference</title></head>
<body>
<h1>NexusIQ REST API v3.2 — Technical Reference</h1>
<p>This document provides a complete reference for the NexusIQ REST API v3.2, including authentication,
rate limiting, endpoint descriptions, request/response schemas, error codes, and code examples.
Base URL: <code>https://api.nexusiq.io/v3</code>. All endpoints require HTTPS.</p>

<h2>Authentication</h2>
<p>NexusIQ supports two authentication mechanisms: Bearer tokens (recommended for server-side integrations)
and API keys (for lightweight scripts and webhooks). All API requests must include an Authorization header.</p>
<pre><code>Authorization: Bearer eyJhbGciOiJSUzI1NiIsInR5cCI6IkpXVCJ9...</code></pre>
<p>Tokens are issued via the OAuth 2.0 client credentials flow. Token lifetime is 3600 seconds.
API keys are issued from the Settings &gt; API Keys console and never expire unless revoked.</p>

<h2>Rate Limiting</h2>
<p>Rate limits are enforced per API key at the account level. The default tier allows 1,000 requests
per minute. Enterprise accounts may request elevated limits. Rate limit headers are included in every response.</p>
<table border="1" cellpadding="4">
  <tr><th>Header</th><th>Description</th><th>Example</th></tr>
  <tr><td>X-RateLimit-Limit</td><td>Max requests per window</td><td>1000</td></tr>
  <tr><td>X-RateLimit-Remaining</td><td>Remaining requests in current window</td><td>842</td></tr>
  <tr><td>X-RateLimit-Reset</td><td>Unix timestamp when limit resets</td><td>1748361600</td></tr>
  <tr><td>Retry-After</td><td>Seconds until retry (on 429 only)</td><td>17</td></tr>
</table>

<h2>Endpoints</h2>

<h3>Documents</h3>

<h4>POST /documents</h4>
<p>Ingest a new document into the pipeline. Accepts multipart/form-data with the file and metadata.
Returns a job_id for tracking pipeline progress via WebSocket or polling.</p>
<pre><code>curl -X POST https://api.nexusiq.io/v3/documents \
  -H "Authorization: Bearer TOKEN" \
  -F "file=@report.pdf" \
  -F "pipeline=custom" \
  -F "metadata={\"source\":\"upload\",\"tags\":[\"medical\",\"trial\"]}"</code></pre>
<p>Response (202 Accepted):</p>
<pre><code>{
  "job_id": "j_01HXYZ1234",
  "status": "queued",
  "pipeline": "custom",
  "created_at": "2026-05-17T10:00:00Z",
  "estimated_duration_s": 45
}</code></pre>

<h4>GET /documents/{doc_id}</h4>
<p>Retrieve metadata and extraction results for a previously ingested document.</p>
<pre><code>curl -X GET https://api.nexusiq.io/v3/documents/doc_abc123 \
  -H "Authorization: Bearer TOKEN"</code></pre>

<h4>GET /documents</h4>
<p>List all documents with optional filters. Supports pagination via cursor-based navigation.</p>
<table border="1" cellpadding="4">
  <tr><th>Parameter</th><th>Type</th><th>Required</th><th>Description</th></tr>
  <tr><td>limit</td><td>integer</td><td>No</td><td>Results per page (default 20, max 100)</td></tr>
  <tr><td>cursor</td><td>string</td><td>No</td><td>Pagination cursor from previous response</td></tr>
  <tr><td>pipeline</td><td>string</td><td>No</td><td>Filter by pipeline: custom | docling</td></tr>
  <tr><td>status</td><td>string</td><td>No</td><td>Filter by status: queued | running | completed | error</td></tr>
  <tr><td>from</td><td>ISO8601</td><td>No</td><td>Filter by created_at &gt;= from</td></tr>
  <tr><td>to</td><td>ISO8601</td><td>No</td><td>Filter by created_at &lt;= to</td></tr>
</table>

<h4>DELETE /documents/{doc_id}</h4>
<p>Delete a document and all associated chunks, vectors, and metadata from the store.
This operation is irreversible. Returns 204 No Content on success.</p>

<h3>Search / Retrieval</h3>

<h4>POST /search</h4>
<p>Perform hybrid (dense + sparse) retrieval over the indexed document corpus.
Supports dense vector search (ANN), BM25 sparse retrieval, and RRF fusion re-ranking.</p>
<pre><code>curl -X POST https://api.nexusiq.io/v3/search \
  -H "Authorization: Bearer TOKEN" \
  -H "Content-Type: application/json" \
  -d '{
    "query": "What was the primary endpoint result for XR-2041?",
    "top_k": 10,
    "pipeline": "custom",
    "filters": {"doc_type": "research_paper", "domain": "medical"},
    "rerank": true,
    "include_metadata": true
  }'</code></pre>

<h3>Jobs</h3>

<h4>GET /jobs/{job_id}</h4>
<p>Poll job status. For real-time updates, use the WebSocket endpoint instead.</p>

<h4>WebSocket ws://api.nexusiq.io/v3/ws/{job_id}</h4>
<p>Connect to receive real-time stage events as the pipeline processes a document.
Each message is a JSON object conforming to the StageEvent schema.</p>

<h2>Error Codes</h2>
<table border="1" cellpadding="4">
  <tr><th>HTTP Status</th><th>Code</th><th>Description</th></tr>
  <tr><td>400</td><td>INVALID_REQUEST</td><td>Malformed request body or missing required field</td></tr>
  <tr><td>401</td><td>UNAUTHORIZED</td><td>Missing or invalid authentication credential</td></tr>
  <tr><td>403</td><td>FORBIDDEN</td><td>Authenticated but insufficient permissions for resource</td></tr>
  <tr><td>404</td><td>NOT_FOUND</td><td>Requested resource does not exist</td></tr>
  <tr><td>409</td><td>CONFLICT</td><td>Resource already exists or concurrent modification</td></tr>
  <tr><td>413</td><td>PAYLOAD_TOO_LARGE</td><td>File exceeds 50 MB limit</td></tr>
  <tr><td>415</td><td>UNSUPPORTED_MEDIA_TYPE</td><td>File type not supported by requested pipeline</td></tr>
  <tr><td>422</td><td>UNPROCESSABLE_ENTITY</td><td>Semantically invalid request (e.g. bad date range)</td></tr>
  <tr><td>429</td><td>RATE_LIMITED</td><td>Too many requests — see Retry-After header</td></tr>
  <tr><td>500</td><td>INTERNAL_ERROR</td><td>Unexpected server error — contact support</td></tr>
  <tr><td>503</td><td>SERVICE_UNAVAILABLE</td><td>Planned maintenance or temporary overload</td></tr>
</table>

<h2>SDK Examples</h2>
<h3>Python</h3>
<pre><code>from nexusiq import NexusIQ
client = NexusIQ(api_key="nxiq_live_xxxxx")
job = client.documents.ingest("report.pdf", pipeline="custom")
result = job.wait()  # blocks until pipeline completes
chunks = result.chunks
print(f"Extracted {len(chunks)} chunks from {result.page_count} pages")</code></pre>

<h3>Node.js</h3>
<pre><code>import { NexusIQ } from '@nexusiq/sdk';
const client = new NexusIQ({ apiKey: process.env.NEXUSIQ_KEY });
const job = await client.documents.ingest(fs.createReadStream('report.pdf'));
const result = await job.wait();
console.log(`${result.chunks.length} chunks, ${result.wordCount} words`);</code></pre>
</body>
</html>""".encode()


# ── 05 Board presentation PPTX ─────────────────────────────────────────────────

def make_board_pptx() -> bytes:
    prs = pptx.Presentation()
    blank = prs.slide_layouts[6]

    slides_content = [
        ("Series B Fundraising — Board Update Q2 2026",
         ["Nexus Software Inc.  |  Confidential",
          "May 2026  |  Board of Directors Meeting",
          "Presenter: Priya Sharma, CEO"]),
        ("Agenda",
         ["1. Q2 Business Highlights",
          "2. ARR & Revenue Performance",
          "3. Product Milestones",
          "4. GTM Traction & Pipeline",
          "5. Series B Strategy & Use of Funds",
          "6. 18-Month Financial Plan",
          "7. Key Risks & Mitigations",
          "8. Board Asks"]),
        ("Q2 2026 Highlights — Record Quarter",
         ["ARR: $47.1M (+138% YoY)  |  Net New ARR: $8.4M (best quarter ever)",
          "NRR: 118%  |  Gross Margin: 74%  |  CAC Payback: 14 months",
          "New Logos: 24 (incl. 3 Fortune 500 accounts — GE Healthcare, Citi, Boeing)",
          "Product: v3.2 GA shipped; SOC 2 Type II renewed; ISO 27001 in progress",
          "Team: 89 FTEs (from 61 in Q4-25); VP Engineering and CRO hired",
          "Pipeline: $32M qualified pipeline; 8 deals >$500K ACV in final stages"]),
        ("ARR Growth Trajectory",
         ["ARR ($ millions):",
          "  Jan-25: $18.4M    |    Jun-25: $29.8M    |    Dec-25: $38.7M",
          "  Mar-26: $43.2M    |    Jun-26: $47.1M (projected: $50.0M)",
          "",
          "Net Revenue Retention by cohort:",
          "  2023 cohort: 124%  |  2024 cohort: 119%  |  2025 cohort: 116% (annualised)",
          "",
          "Revenue Mix: 82% subscription / 18% professional services",
          "Enterprise (>$100K ACV): 48% of ARR (up from 31% in Q2-25)"]),
        ("Product — v3.2 & Roadmap",
         ["v3.2 Released April 2026:",
          "  - Hybrid RAG pipeline with BM25 + ANN fusion (latency: 87ms p95)",
          "  - Multi-modal ingestion: PDF, DOCX, XLSX, PPTX, HTML",
          "  - LLM-graded semantic verification (L2) with Claude integration",
          "  - Enterprise SSO: Okta, Azure AD, Google Workspace",
          "  - HIPAA BAA now available; FedRAMP Moderate in assessment",
          "",
          "H2 2026 Roadmap (board approval requested):",
          "  - Real-time collaborative annotations (Q3)",
          "  - Graph-based retrieval (Q3)",
          "  - On-premises deployment option for air-gapped customers (Q4)",
          "  - Multi-tenancy isolation improvements (Q4)"]),
        ("GTM — Pipeline & Expansion",
         ["Sales Pipeline (as of May 17, 2026):",
          "  Stage 1 Discovery: 142 opps ($18.2M)   |  Stage 2 Demo: 67 ($11.4M)",
          "  Stage 3 Proposal: 31 ($7.8M)   |   Stage 4 Negotiation: 8 ($4.6M)",
          "  Weighted Pipeline: $22.3M  |  Coverage ratio: 3.2x vs $7M target",
          "",
          "Key Wins Q2:",
          "  - GE Healthcare: $1.2M ACV (3-yr) — clinical document intelligence",
          "  - Citi: $840K ACV — regulatory filing automation",
          "  - Boeing: $620K ACV — technical manual search",
          "",
          "Expansion Highlights: 14 upsell/cross-sell deals totalling $2.1M net new ARR"]),
        ("Series B — Strategy & Terms",
         ["Target Raise: $50M Series B",
          "Valuation: $310M pre-money (6.6x ARR at $47.1M)",
          "Lead Investor: Andreessen Horowitz (term sheet signed)",
          "Follow-on: Existing investors Sequoia and Founders Fund (pro-rata)",
          "",
          "Use of Funds:",
          "  40% Product & Engineering (30 hires, infra, security certs)",
          "  35% Sales & Marketing (20 AE hires, demand gen, events)",
          "  15% Customer Success & Professional Services",
          "  10% G&A, Legal, Working Capital",
          "",
          "Runway: 22 months to Series C or profitability gate at $85M ARR"]),
        ("18-Month Financial Plan",
         ["Revenue Plan (base case):",
          "  Q3-26: $13.9M  |  Q4-26: $16.2M  |  FY-26: $53.4M",
          "  Q1-27: $18.8M  |  Q2-27: $21.9M  |  H1-27: $40.7M",
          "",
          "Key Assumptions:",
          "  Monthly growth rate: 12% (consistent with H1-26 actuals)",
          "  NRR: 118% (conservative vs 121% trailing 12-month)",
          "  Gross Margin expansion: 74% → 77% (infra optimisations)",
          "  Headcount: 89 → 138 FTEs by Dec-26 (38 net adds)",
          "",
          "Path to Profitability: EBITDA breakeven at $85M ARR (Q3-27)"]),
        ("Key Risks & Mitigations",
         ["Risk 1: Enterprise sales cycle lengthening",
          "  Mitigation: PLG motion for SMB; dedicated enterprise pods; RevOps tooling",
          "",
          "Risk 2: LLM API cost volatility (Claude, OpenAI)",
          "  Mitigation: Model agnosticism; local inference option for high-volume customers",
          "",
          "Risk 3: Data privacy regulation (EU AI Act, CCPA 2.0)",
          "  Mitigation: Privacy-by-design architecture; on-prem option; DPA templates ready",
          "",
          "Risk 4: Competitive pressure (Microsoft Copilot, Google NotebookLM)",
          "  Mitigation: Domain-specific fine-tuning; enterprise trust/compliance moat",
          "",
          "Risk 5: Key-person dependency (CTO, CRO new hires)",
          "  Mitigation: Equity retention packages; succession planning underway"]),
        ("Board Asks",
         ["1. APPROVE Series B term sheet at $310M pre-money with a16z lead",
          "2. RATIFY use of proceeds allocation (40/35/15/10 split)",
          "3. APPROVE FY-26 budget amendment (+$4.2M for accelerated GTM hiring)",
          "4. ELECT new board member: Partner from a16z (seat approved in Series B term)",
          "5. REVIEW & NOTE: FedRAMP Moderate assessment timeline (target: Q1-27)",
          "",
          "Next Board Meeting: August 19, 2026",
          "Questions & Discussion"]),
    ]

    for title, bullets in slides_content:
        slide = prs.slides.add_slide(blank)
        txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(9.2), Inches(0.7))
        tf = txBox.text_frame
        tf.text = title
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.size = PptPt(22)

        body = slide.shapes.add_textbox(Inches(0.4), Inches(1.1), Inches(9.2), Inches(5.8))
        btf = body.text_frame
        btf.word_wrap = True
        for i, line in enumerate(bullets):
            if i == 0:
                btf.text = line
            else:
                p = btf.add_paragraph()
                p.text = line
            para = btf.paragraphs[i]
            if para.runs:
                para.runs[0].font.size = PptPt(13) if i > 0 else PptPt(14)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    docs = {
        "01_pharmaceutical_trial.pdf": make_pharma_pdf(),
        "02_financial_model.xlsx":     make_financial_xlsx(),
        "03_vendor_contract.docx":     make_contract_docx(),
        "04_technical_spec.html":      make_technical_html(),
        "05_board_presentation.pptx":  make_board_pptx(),
    }
    for filename, data in docs.items():
        path = DEMO_DOCS_DIR / filename
        path.write_bytes(data)
        print(f"  {path.name}  ({len(data):,} bytes)")
    print(f"\nAll demo docs written to {DEMO_DOCS_DIR}/")


if __name__ == "__main__":
    main()
