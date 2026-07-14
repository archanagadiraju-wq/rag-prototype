from __future__ import annotations
import asyncio
import json
import os
import tempfile
import time
import uuid
from pathlib import Path

import structlog
from fastapi import FastAPI, UploadFile, File, Form, WebSocket, WebSocketDisconnect, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import Response, FileResponse
from pydantic import BaseModel

from config import settings
from pipelines.mock import run_mock_pipeline
from pipelines.custom.runner import run_custom_pipeline
from pipelines.docling.runner import run_docling_pipeline
from agent.runner import run_agent_pipeline

log = structlog.get_logger()

DEMO_DOCS_DIR = Path(__file__).parent.parent / "demo_docs"

from api_schemas import (
    JobCreateResponse,
    JobStatusResponse,
    AskRequest,
    AskResponse,
    DemoDocList,
    ErrorResponse,
    StorageSummary,
    JobSummary,
    InspectChunksResponse,
    InspectSqlResponse,
    InspectKgResponse,
)

OPENAPI_DESCRIPTION = """
# RAG Ingestion Engine API

Backend powering the **Agentic RAG** prototype. Provides four document-ingestion
pipelines (Mode A custom, Mode B Docling, Mode C compare-both, Mode D
AI-agent-orchestrated) plus a hybrid retrieval + LLM-grounded Q&A endpoint.

## Workflow

1. **Create a job** with `POST /api/jobs` — upload a file OR pick a demo doc.
2. **Connect to the WebSocket** `/ws/{job_id}` to stream per-stage progress
   events as they fire.
3. **Wait for status='completed'** on `GET /api/jobs/{job_id}` (or for the final
   stage event over the WS).
4. **Ask questions** with `POST /api/jobs/{job_id}/ask` — any number of
   ad-hoc questions against the ingested document.

## Pipelines

| Mode | Stages | Use when |
|------|--------|----------|
| `custom`  | 11 hand-built stages | Born-digital PDFs, office formats. Fast. |
| `docling` | 8 stages (Docling collapses parse/intel/chunk into one) | Scanned PDFs needing OCR, complex multi-page tables. Slower (CPU-bound ML). |
| `compare` | Both A and B concurrently | Side-by-side quality + cost comparison. |
| `agent`   | Dynamic (AI agent picks tools) | Mixed-content docs where the right pipeline isn't obvious. |

## Quality guarantees

- Every Q&A answer is **grounded** in retrieved chunks (no hallucination).
- Every answer is **judged** by an independent Claude call for verdict + score.
- Every stage has L1 **verification checks**; pass rates are visible per stage.
- Every step is **resumable** — backend restart preserves all work via disk cache.
- Every LLM call uses **retry-with-backoff** for 429s and transient 5xx.

## Authentication

Currently none — this is a prototype. Add an auth header check in front of
the FastAPI app before exposing to untrusted networks.
"""

OPENAPI_TAGS = [
    {"name": "Jobs",       "description": "Create and inspect ingestion jobs."},
    {"name": "Q&A",        "description": "Ad-hoc questions against ingested documents."},
    {"name": "Files",      "description": "Download or render the source document."},
    {"name": "Demo Docs",  "description": "Bundled demonstration documents."},
    {"name": "WebSocket",  "description": "Real-time stage progress streaming."},
]

app = FastAPI(
    title="RAG Ingestion Engine",
    version="2.0.0",
    description=OPENAPI_DESCRIPTION,
    openapi_tags=OPENAPI_TAGS,
    contact={"name": "RAG Prototype Team"},
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173", "http://127.0.0.1:5173"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# In-memory job store (Layer 1 — replaced with Redis in later layers)
_jobs: dict[str, dict] = {}
_ws_connections: dict[str, list[WebSocket]] = {}
# Event replay buffer: events emitted before WS connects are stored here
_job_events: dict[str, list[dict]] = {}
# File bytes stored per job for preview endpoint
_job_files: dict[str, dict] = {}  # job_id -> {bytes, filename, mime}
# Track the asyncio.Task running each pipeline so DELETE /api/jobs/{id} can cancel.
_job_tasks: dict[str, asyncio.Task] = {}
# Cooperative cancel signal — set by DELETE, checked by the agent runner between
# turns. Cancelling the asyncio Task raises CancelledError at the next await, but
# the agent loop also reads this so partial work persists cleanly.
_cancelled_jobs: set[str] = set()


_EXT_MIME: dict[str, str] = {
    ".pdf":  "application/pdf",
    ".html": "text/html",
    ".htm":  "text/html",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".txt":  "text/plain",
}

def _mime_for_ext(ext: str) -> str:
    return _EXT_MIME.get(ext, "application/octet-stream")


# ── REST endpoints ─────────────────────────────────────────────────────────────


@app.post(
    "/api/jobs",
    response_model=JobCreateResponse,
    tags=["Jobs"],
    summary="Create a new ingestion job",
    responses={
        200: {"description": "Job created and queued for ingestion."},
        400: {"description": "Missing file/demo_doc, or invalid pipeline.", "model": ErrorResponse},
    },
)
async def create_job(
    file: UploadFile | None = File(None, description="Document to ingest. Supported: PDF, DOCX, PPTX, XLSX, HTML."),
    demo_doc: str | None = Form(None, description="Filename of a bundled demo document. Mutually exclusive with `file`."),
    pipeline: str = Form("custom", description="Pipeline mode: `custom`, `docling`, `compare`, or `agent`."),
):
    """Create a new document-ingestion job and start the selected pipeline.

    Provide **either** a `file` upload **or** a `demo_doc` name (not both). The
    pipeline runs in the background; subscribe to `/ws/{job_id}` to stream
    per-stage progress events, or poll `GET /api/jobs/{job_id}` for status.

    For Mode C (`compare`) both Mode A and Mode B run concurrently against the
    same source document; events for each are tagged with their `pipeline` field.

    For Mode D (`agent`) an AI agent picks the right tools dynamically — stage
    IDs are not fixed; stages are emitted as they're discovered.
    """
    job_id = str(uuid.uuid4())
    created_at = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())

    filepath: Path | None = None
    source_type = "demo_doc"
    tmp_path: str | None = None

    if file and file.filename:
        suffix = Path(file.filename).suffix or ".bin"
        contents = await file.read()
        fd, tmp_path = tempfile.mkstemp(suffix=suffix)
        with os.fdopen(fd, "wb") as fh:
            fh.write(contents)
        filepath = Path(tmp_path)
        source_type = "upload"
        _job_files[job_id] = {
            "bytes": contents,
            "filename": file.filename,
            "mime": file.content_type or "application/octet-stream",
        }
    elif demo_doc:
        filepath = DEMO_DOCS_DIR / demo_doc
        source_type = "demo_doc"
        if filepath.exists():
            _job_files[job_id] = {
                "bytes": filepath.read_bytes(),
                "filename": demo_doc,
                "mime": _mime_for_ext(filepath.suffix.lower()),
            }

    filename = (filepath.name if filepath else None) or "unknown"
    _jobs[job_id] = {
        "job_id": job_id,
        "status": "running",
        "pipeline": pipeline,
        "doc_filename": filename,
        "created_at": created_at,
        "completed_at": None,
        "error": None,
    }
    _ws_connections[job_id] = []
    _job_events[job_id] = []

    task = asyncio.create_task(_run_pipeline(job_id, pipeline, filepath, source_type, tmp_path))
    _job_tasks[job_id] = task

    return JobCreateResponse(job_id=job_id, created_at=created_at)


@app.delete(
    "/api/jobs/{job_id}",
    tags=["Jobs"],
    summary="Cancel a running job",
    responses={
        200: {"description": "Cancel signal sent. Task already completed jobs return 200 too (no-op)."},
        404: {"description": "No job with that ID.", "model": ErrorResponse},
    },
)
async def cancel_job(job_id: str):
    """Cancel a running ingestion job.

    Marks the job as cancelled, sets a cooperative cancel flag the agent
    runner checks between turns, and calls `task.cancel()` to break the
    current `await`. Long synchronous work (Docling parse via
    `asyncio.to_thread`) cannot be killed cleanly — the worker thread keeps
    running until it finishes, but the asyncio coroutine returns immediately
    and the job is marked cancelled, so the user gets instant feedback.
    """
    if job_id not in _jobs:
        raise HTTPException(status_code=404, detail="Job not found")

    if _jobs[job_id].get("status") in ("completed", "cancelled", "error"):
        return {"job_id": job_id, "status": _jobs[job_id]["status"], "noop": True}

    _cancelled_jobs.add(job_id)
    task = _job_tasks.get(job_id)
    if task and not task.done():
        task.cancel()
    return {"job_id": job_id, "status": "cancel_requested", "noop": False}


def is_cancelled(job_id: str) -> bool:
    """Cooperative cancel check — agent runner calls this between turns."""
    return job_id in _cancelled_jobs


@app.get(
    "/api/jobs/{job_id}",
    response_model=JobStatusResponse,
    tags=["Jobs"],
    summary="Get job status",
    responses={
        200: {"description": "Job state returned."},
        404: {"description": "No job with that ID.", "model": ErrorResponse},
    },
)
async def get_job(job_id: str):
    """Return the current high-level state of a job.

    For real-time stage-level updates, subscribe to `/ws/{job_id}` instead.
    This endpoint reports only the rollup status: `queued`, `running`,
    `completed`, or `error` (with `error` populated).
    """
    if job_id not in _jobs:
        raise HTTPException(status_code=404, detail="Job not found")
    return _jobs[job_id]


@app.post(
    "/api/jobs/{job_id}/ask",
    response_model=AskResponse,
    tags=["Q&A"],
    summary="Ask an ad-hoc question",
    responses={
        200: {"description": "Answer with retrieved chunks, judge verdict, and full audit trail."},
        400: {"description": "Empty question.", "model": ErrorResponse},
        404: {"description": "Job not ingested yet (no embedded chunks in cache).", "model": ErrorResponse},
    },
)
async def ask_question(job_id: str, body: AskRequest):
    """Run an ad-hoc retrieval + LLM Q&A against an ingested document.

    Same hybrid retrieval (dense + BM25 + KG) + Claude haiku answer + LLM-as-judge
    pipeline that powers the built-in showcase questions, but with your question.
    The job must have reached the Vector Store stage so `embedded_chunks` are
    present on disk.

    Disk-backed cache → works across backend restarts.

    **Pipeline parameter**: pass `custom` for Mode A or Mode D (agent) — they
    share the same cache prefix. Pass `docling` for Mode B. Other values
    silently fall back to `custom`.

    **Response includes**:
    - The grounded answer
    - The 5 retrieved chunks with scores (for verification)
    - The exact system + user prompt sent to the model (audit trail)
    - Judge verdict + score + rationale (independent quality check)
    - Per-call tokens, cost, latency
    """
    from pipelines.custom import stage_11_llm_answer
    if not body.question or not body.question.strip():
        raise HTTPException(status_code=400, detail="question required")
    pipeline = body.pipeline if body.pipeline in ("custom", "docling") else "custom"
    cache_prefix = "d_" if pipeline == "docling" else ""
    result = await stage_11_llm_answer.answer_one(
        question=body.question.strip(),
        job_id=job_id,
        cache_prefix=cache_prefix,
    )
    if result.get("error"):
        raise HTTPException(
            status_code=404,
            detail=f"{result['error']}. Run the {pipeline} pipeline first.",
        )
    return result


# ── Storage summary ───────────────────────────────────────────────────────────

# Haiku 4.5 + text-embedding-3-large pricing — single source of truth
_PRICE_HAIKU_IN      = 0.80  / 1_000_000
_PRICE_HAIKU_OUT     = 4.00  / 1_000_000
_PRICE_HAIKU_CACHE_R = 0.08  / 1_000_000   # cache read at 10% of normal
_PRICE_HAIKU_CACHE_W = 1.00  / 1_000_000   # cache write at 1.25× normal
_PRICE_OPENAI_EMBED  = 0.13  / 1_000_000


def _build_job_summary(job_id: str) -> dict | None:
    """Walk all persisted stage payloads + the job record and produce a
    consolidated summary: total wall time, tokens by provider, cache effect,
    grand total cost, per-stage breakdown.

    Reads from `services.job_cache._stage_done_*` keys (disk-backed) so this
    works for any job regardless of pipeline (Mode A, B, agent) and survives
    backend restarts.
    """
    import services.job_cache as cache

    job_meta = _jobs.get(job_id) or {}
    job_dir = Path(__file__).resolve().parent / "data" / "jobs" / job_id
    if not job_dir.exists() and not job_meta:
        return None

    # Find all persisted stage_done events on disk
    stage_records: list[dict] = []
    if job_dir.exists():
        for f in sorted(job_dir.glob("_stage_done_*.json")):
            try:
                with open(f) as fh:
                    rec = json.load(fh)
                rec["_filename"] = f.name
                stage_records.append(rec)
            except Exception:
                pass

    # Sort stages by their numeric ID parsed from filename (custom_4, agent_1, etc.)
    def _sid(rec):
        try:
            return int(rec["_filename"].rsplit("_", 1)[-1].replace(".json", ""))
        except Exception:
            return 0
    stage_records.sort(key=_sid)

    # ── Per-stage aggregation ──────────────────────────────────────────────
    stages: list[dict] = []
    anth_input = anth_output = anth_cache_r = anth_cache_w = 0
    oai_embedding_tokens = 0
    iterations: int | None = None

    for rec in stage_records:
        payload = rec.get("payload") or {}
        duration_ms = rec.get("duration_ms", 0) or 0
        stage_id = _sid(rec)
        name = payload.get("stage_name") or rec.get("_filename", "").replace("_stage_done_", "").replace(".json", "")

        # Mode A/B stage payload: top-level `llm_*` fields are Anthropic OR OpenAI
        # depending on stage name (Embedding is OpenAI; everything else Anthropic).
        is_openai_embed_stage = "embedding" in name.lower() and "use_real_embeddings" in payload
        s_in  = int(payload.get("llm_input_tokens") or 0)
        s_out = int(payload.get("llm_output_tokens") or 0)

        # Agent stage payload nests metrics under `turn`
        turn = payload.get("turn") or {}
        if turn:
            s_in  = max(s_in, int(turn.get("turn_input_tokens") or 0))
            s_out = max(s_out, int(turn.get("turn_output_tokens") or 0))
            anth_cache_r += int(turn.get("turn_cache_read_tokens") or 0)
            anth_cache_w += int(turn.get("turn_cache_create_tokens") or 0)

        # "Agent finished" stage carries the rolled-up totals from the runner —
        # use these as authoritative (max() handles cases where the per-stage
        # sum is incomplete, e.g., an old job persisted before the no-tool-turn
        # fix landed).
        if "agent finished" in name.lower() or "iterations" in payload:
            if payload.get("iterations") is not None:
                iterations = payload["iterations"]
            anth_input   = max(anth_input,   int(payload.get("total_input_tokens")   or 0))
            anth_output  = max(anth_output,  int(payload.get("total_output_tokens")  or 0))
            anth_cache_r = max(anth_cache_r, int(payload.get("cache_read_tokens")    or 0))
            anth_cache_w = max(anth_cache_w, int(payload.get("cache_create_tokens")  or 0))

        # Tool execution result — tools can themselves make API calls that
        # don't appear in the agent's per-turn metrics:
        #   • agent.embed_and_index → OpenAI embeddings (embedding_cost_usd)
        #   • agent.describe_tables → Claude table descriptions (llm_input_tokens)
        # Both must be added to provider totals to match raw API ground truth.
        tool_result = payload.get("tool_result") or {}
        tool_anth_in  = int(tool_result.get("llm_input_tokens")  or 0)
        tool_anth_out = int(tool_result.get("llm_output_tokens") or 0)
        oai_real_cost_this_stage = float(tool_result.get("embedding_cost_usd") or 0.0)

        # Route to provider buckets
        if is_openai_embed_stage:
            oai_embedding_tokens += s_in
        else:
            anth_input  += s_in
            anth_output += s_out

        # Tool-internal Anthropic calls (e.g., stage_06_multimodal descriptions
        # invoked by agent.describe_tables / auto.describe_tables). These DON'T
        # go through the cached prefix so they contribute to uncached input.
        anth_input  += tool_anth_in
        anth_output += tool_anth_out

        # Compute per-stage cost. Agent stages know their cost via turn.turn_cost_usd
        if turn and "turn_cost_usd" in turn:
            stage_cost = float(turn["turn_cost_usd"])
        elif is_openai_embed_stage:
            stage_cost = s_in * _PRICE_OPENAI_EMBED
        else:
            stage_cost = s_in * _PRICE_HAIKU_IN + s_out * _PRICE_HAIKU_OUT

        stages.append({
            "stage_id":      stage_id,
            "name":          name,
            "duration_ms":   round(duration_ms, 1),
            "input_tokens":  s_in,
            "output_tokens": s_out,
            "cost_usd":      round(stage_cost, 6),
        })

    # ── Cost math ───────────────────────────────────────────────────────────
    anth_cost = (
        anth_input    * _PRICE_HAIKU_IN
        + anth_output * _PRICE_HAIKU_OUT
        + anth_cache_r * _PRICE_HAIKU_CACHE_R
        + anth_cache_w * _PRICE_HAIKU_CACHE_W
    )
    # No-cache baseline: every cached_read token would have been a normal input
    no_cache_baseline = (
        (anth_input + anth_cache_r + anth_cache_w) * _PRICE_HAIKU_IN
        + anth_output * _PRICE_HAIKU_OUT
    )
    saved = max(0.0, no_cache_baseline - anth_cost)
    saved_pct = (saved / no_cache_baseline * 100) if no_cache_baseline > 0 else 0.0

    # Prefer the tool-reported actual OpenAI cost (sum across stages); fall back
    # to the estimate-from-tokens path if no tool result carried it.
    oai_real_cost_total = sum(
        float((rec.get("payload", {}).get("tool_result") or {}).get("embedding_cost_usd") or 0.0)
        for rec in stage_records
    )
    if oai_real_cost_total > 0:
        oai_cost = oai_real_cost_total
        # Approximate token count from cost since the API reported dollars not tokens
        oai_embedding_tokens = max(
            oai_embedding_tokens,
            int(oai_real_cost_total / _PRICE_OPENAI_EMBED),
        )
    else:
        oai_cost = oai_embedding_tokens * _PRICE_OPENAI_EMBED

    # ── Wall time ───────────────────────────────────────────────────────────
    wall_time_s = 0.0
    if job_meta.get("created_at") and job_meta.get("completed_at"):
        try:
            from datetime import datetime
            t0 = datetime.strptime(job_meta["created_at"], "%Y-%m-%dT%H:%M:%SZ")
            t1 = datetime.strptime(job_meta["completed_at"], "%Y-%m-%dT%H:%M:%SZ")
            wall_time_s = (t1 - t0).total_seconds()
        except Exception:
            pass
    if wall_time_s == 0.0:
        # Fall back: sum of stage durations
        wall_time_s = sum(s["duration_ms"] for s in stages) / 1000.0

    return {
        "job_id":     job_id,
        "pipeline":   job_meta.get("pipeline"),
        "filename":   job_meta.get("doc_filename"),
        "status":     job_meta.get("status", "unknown"),
        "wall_time_s": round(wall_time_s, 2),
        "iterations": iterations,
        "anthropic": {
            "input_tokens":         anth_input,
            "output_tokens":        anth_output,
            "cache_create_tokens":  anth_cache_w,
            "cache_read_tokens":    anth_cache_r,
            "cost_usd":             round(anth_cost, 6),
            "no_cache_baseline_usd": round(no_cache_baseline, 6),
            "saved_usd":            round(saved, 6),
            "saved_pct":            round(saved_pct, 1),
        },
        "openai": {
            "embedding_tokens": oai_embedding_tokens,
            "cost_usd":         round(oai_cost, 6),
        },
        "total_cost_usd": round(anth_cost + oai_cost, 6),
        "total_tokens":   anth_input + anth_output + anth_cache_r + anth_cache_w + oai_embedding_tokens,
        "stages":         stages,
    }


@app.get(
    "/api/jobs/{job_id}/summary",
    response_model=JobSummary,
    tags=["Jobs"],
    summary="Consolidated cost/token/time summary for one job",
    responses={
        200: {"description": "Job totals: wall time, tokens by provider, cache effect, grand cost, per-stage breakdown."},
        404: {"description": "Job not found.", "model": ErrorResponse},
    },
)
async def get_job_summary(job_id: str):
    """Aggregate everything spent on one document — across both LLM providers.

    Walks every persisted stage payload for the job and computes:
    - **Wall time** (creation → completion)
    - **Anthropic** totals: input/output/cache_read/cache_write tokens, cost,
      no-cache baseline, savings %
    - **OpenAI** totals: embedding tokens + cost
    - **Grand total** = Anthropic + OpenAI
    - **Per-stage** breakdown in execution order

    Works for any pipeline (Mode A, B, C, D) — sources stage data from disk,
    so survives backend restarts.
    """
    summary = _build_job_summary(job_id)
    if summary is None:
        raise HTTPException(status_code=404, detail="Job not found")
    return summary


def _survey_storage(job_id: str) -> dict:
    """Inspect every storage sink for a job_id and return a consolidated dict.

    Reads from: services.job_cache (memory + disk fallback), the filesystem
    job directory, and SQLite if present.
    """
    import sqlite3
    import services.job_cache as cache
    job_dir = Path(__file__).resolve().parent / "data" / "jobs" / job_id

    # Try both un-prefixed (Mode A / agent) and d_-prefixed (Mode B) keys
    # so the same endpoint works regardless of which pipeline ingested.
    def _get_either(key: str, default=None):
        v = cache.get(job_id, key)
        if v is None:
            v = cache.get(job_id, f"d_{key}")
        return v if v is not None else default

    # 1. Qdrant
    qdrant_collection = _get_either("qdrant_collection")
    embedded = _get_either("embedded_chunks", []) or []
    sample_vector: list[float] = []
    if embedded:
        first_vec = (embedded[0] or {}).get("vector") or []
        if first_vec:
            sample_vector = [round(float(x), 5) for x in first_vec[:8]]

    # Liveness probe against the real Qdrant
    qdrant_live = False
    if qdrant_collection:
        try:
            from qdrant_client import QdrantClient
            QdrantClient(host="localhost", port=6333, timeout=2).get_collections()
            qdrant_live = True
        except Exception:
            qdrant_live = False

    qdrant = {
        "collection":         qdrant_collection if isinstance(qdrant_collection, str) else None,
        "vectors":            len(embedded),
        "dimensions":         1536,
        "live":               qdrant_live,
        "distance":           "COSINE",
        "hnsw_m":             8,
        "hnsw_ef_construct":  100,
        "embedding_model":    "text-embedding-3-large",
        "sample_vector":      sample_vector,
    }

    # 2. SQLite — look for tables.db / tables_.db / tables_d_.db in the job dir
    sqlite_data = {"file": None, "size_bytes": 0, "tables": []}
    if job_dir.exists():
        for candidate in ("tables_.db", "tables.db", "tables_d_.db"):
            db_path = job_dir / candidate
            if not db_path.exists():
                continue
            sqlite_data["file"]       = candidate
            sqlite_data["size_bytes"] = db_path.stat().st_size
            try:
                conn = sqlite3.connect(str(db_path))
                table_rows = conn.execute(
                    "SELECT name FROM sqlite_master WHERE type='table'"
                ).fetchall()
                for (tname,) in table_rows:
                    n = conn.execute(f'SELECT COUNT(*) FROM "{tname}"').fetchone()[0]
                    cols = [r[1] for r in conn.execute(f'PRAGMA table_info("{tname}")').fetchall()]
                    sqlite_data["tables"].append({"name": tname, "rows": n, "columns": cols})
                conn.close()
            except Exception as exc:
                log.warning("storage_survey_sqlite", job_id=job_id[:8], error=str(exc))
            break

    # 3. BM25
    bm25_idx = _get_either("bm25_index") or {}
    bm25 = {
        "unique_terms": len(bm25_idx.get("idf", {})) if isinstance(bm25_idx, dict) else 0,
        "doc_count":    bm25_idx.get("N", 0) if isinstance(bm25_idx, dict) else 0,
        "avg_doc_len":  round(bm25_idx.get("avgdl", 0.0) if isinstance(bm25_idx, dict) else 0.0, 1),
    }

    # 4. Knowledge graph
    kg = _get_either("knowledge_graph")
    kg_info = {"nodes": 0, "edges": 0, "entity_types": []}
    if kg is not None and hasattr(kg, "number_of_nodes"):
        kg_info["nodes"] = kg.number_of_nodes()
        kg_info["edges"] = kg.number_of_edges()
        types: set[str] = set()
        for _node, data in kg.nodes(data=True):
            label = data.get("label")
            if label:
                types.add(label)
        kg_info["entity_types"] = sorted(types)

    # 5. Disk
    disk_files: list[dict] = []
    disk_total = 0
    if job_dir.exists():
        for f in sorted(job_dir.iterdir()):
            if not f.is_file():
                continue
            sz = f.stat().st_size
            disk_total += sz
            disk_files.append({"name": f.name, "size_bytes": sz})
    disk_info = {
        "path":        f"backend/data/jobs/{job_id}/",
        "total_bytes": disk_total,
        "file_count":  len(disk_files),
        "files":       disk_files,
    }

    # 6. In-memory cache (the live store dict)
    mem_keys = list(cache._store.get(job_id, {}).keys())
    cache_info = {"key_count": len(mem_keys), "keys": sorted(mem_keys)}

    # 7. Document facts (typed, single-valued document properties; per-doc
    #    JSON. See services/fact_extractor.py — separate from SQL line-item
    #    tables.)
    from services.fact_extractor import load_facts
    facts_payload = load_facts(job_id) or load_facts(job_id, cache_prefix="d_")
    if facts_payload:
        kinds: dict[str, int] = {}
        for f in facts_payload.get("facts", []) or []:
            t = f.get("type", "text")
            kinds[t] = kinds.get(t, 0) + 1
        facts_info = {
            "extracted":         True,
            "fact_count":        len(facts_payload.get("facts", []) or []),
            "by_type":           kinds,
            "extractor_version": facts_payload.get("extractor_version"),
            "extracted_at":      facts_payload.get("extracted_at"),
            "stats":             facts_payload.get("stats") or {},
            "error":             facts_payload.get("error"),
        }
    else:
        facts_info = {
            "extracted":         False,
            "fact_count":        0,
            "by_type":           {},
            "extractor_version": None,
            "extracted_at":      None,
            "stats":             {},
            "error":             None,
        }

    return {
        "job_id": job_id,
        "exists": bool(job_dir.exists() or mem_keys),
        "qdrant":  qdrant,
        "sqlite":  sqlite_data,
        "bm25":    bm25,
        "kg":      kg_info,
        "facts":   facts_info,
        "disk":    disk_info,
        "cache":   cache_info,
    }


@app.get(
    "/api/jobs/{job_id}/storage",
    response_model=StorageSummary,
    tags=["Jobs"],
    summary="Consolidated storage summary for one job",
    responses={
        200: {"description": "All storage sinks (Qdrant, SQLite, BM25, KG, disk, in-memory) for this job."},
        404: {"description": "No storage state found for this job_id.", "model": ErrorResponse},
    },
)
async def get_job_storage(job_id: str):
    """Return a consolidated view of every storage format the document was saved into.

    Surveys six sinks for the given job:

    1. **Qdrant** — vector collection name + vector count
    2. **SQLite** — structured tables extracted from the document
    3. **BM25** — sparse keyword index size
    4. **Knowledge graph** — entity + relationship counts
    5. **Disk** — every file written under `backend/data/jobs/<job_id>/`
    6. **In-memory cache** — currently-hot keys

    Returns 404 if the job has no on-disk or in-memory state — likely because
    it was never created or the data was cleared.
    """
    survey = _survey_storage(job_id)
    if not survey["exists"]:
        raise HTTPException(status_code=404, detail="No storage state for this job")
    return survey


# ── Document facts endpoints ─────────────────────────────────────────────────


@app.get(
    "/api/jobs/{job_id}/facts",
    tags=["Jobs"],
    summary="All extracted document facts (typed JSON property store)",
    responses={
        200: {"description": "facts.json payload."},
        404: {"description": "Job has no facts extracted yet.", "model": ErrorResponse},
    },
)
async def get_facts(job_id: str):
    """Return the per-document facts.json: typed single-valued properties
    (capacity, budget, dates, approver, ...) with source citation."""
    from services.fact_extractor import load_facts
    payload = load_facts(job_id) or load_facts(job_id, cache_prefix="d_")
    if payload is None:
        raise HTTPException(status_code=404, detail="No facts extracted for this job")
    return payload


@app.post(
    "/api/jobs/{job_id}/columns/describe",
    tags=["Jobs"],
    summary="Generate semantic column descriptions for every SQL table",
    responses={
        200: {"description": "Descriptions added to sql_registry; report returned."},
        404: {"description": "Job has no SQL registry.", "model": ErrorResponse},
    },
)
async def describe_columns_endpoint(job_id: str):
    """Run a Claude pass per SQL table to write one-line semantic descriptions
    for every column, plus a scope tag (this_project / reference_data / ...).

    The descriptions get injected into the SQL router's prompt so it can
    disambiguate columns with similar names across different tables — the
    root cause of the SQL routing regression we measured in the eval.
    """
    import services.job_cache as cache
    from services.column_describer import describe_columns

    cache_prefix = ""
    if not cache.get(job_id, "sql_registry"):
        cache_prefix = "d_"
        if not cache.get(job_id, "d_sql_registry"):
            raise HTTPException(status_code=404, detail="No SQL registry for this job")

    report = await describe_columns(job_id, cache_prefix=cache_prefix)
    return report


@app.post(
    "/api/jobs/{job_id}/facts/extract",
    tags=["Jobs"],
    summary="Run (or re-run) fact extraction for an ingested job",
    responses={
        200: {"description": "Extraction completed; returns the new facts.json payload."},
        404: {"description": "Job has no embedded chunks (not ingested yet).", "model": ErrorResponse},
    },
)
async def extract_facts_endpoint(job_id: str):
    """Synchronously extract single-valued document properties via Claude
    and persist as `facts.json` under `data/jobs/{job_id}/`.

    Idempotent — running again overwrites the previous facts.json. The output
    contains a `stats` block (extracted vs rejected counts, tokens, latency)
    so callers can monitor extraction quality without parsing logs.
    """
    import services.job_cache as cache
    from services.fact_extractor import extract_facts

    has_chunks = bool(
        cache.get(job_id, "embedded_chunks") or cache.get(job_id, "d_embedded_chunks")
    )
    if not has_chunks:
        # Try to hydrate from disk
        job_dir = Path(__file__).resolve().parent / "data" / "jobs" / job_id
        if not (job_dir / "embedded_chunks.json").exists():
            raise HTTPException(status_code=404, detail="No embedded chunks for this job")

    # Pick the cache_prefix that actually has data
    cache_prefix = "" if cache.get(job_id, "embedded_chunks") else "d_"
    payload = await extract_facts(job_id, cache_prefix=cache_prefix)
    return payload


# ── Storage inspect endpoints — full drill-down for verification ────────────


def _embedded_chunks_either(job_id: str) -> list[dict]:
    """Return embedded_chunks regardless of pipeline (un-prefixed or d_ for Mode B)."""
    import services.job_cache as cache
    v = cache.get(job_id, "embedded_chunks", []) or []
    if not v:
        v = cache.get(job_id, "d_embedded_chunks", []) or []
    return v


@app.get(
    "/api/jobs/{job_id}/inspect/chunks",
    response_model=InspectChunksResponse,
    tags=["Jobs"],
    summary="Full paginated list of embedded chunks for one job",
    responses={
        200: {"description": "Paginated chunks in the order they were embedded."},
        404: {"description": "No chunks for this job.", "model": ErrorResponse},
    },
)
async def inspect_chunks(
    job_id: str,
    offset: int = 0,
    limit: int = 50,
    chunk_type: str | None = None,
):
    """Stream the full content of every chunk that landed in the vector DB.

    Use `chunk_type` to filter (e.g. `table_summary`, `prose`, `ocr_prose`).
    Each item includes the first 8 dims of its stored vector so you can
    verify real embeddings (not mocks) landed.
    """
    limit = max(1, min(limit, 500))
    offset = max(0, offset)

    chunks = _embedded_chunks_either(job_id)
    if not chunks:
        raise HTTPException(status_code=404, detail="No chunks for this job")

    distinct_types = sorted({
        ((c.get("metadata") or {}).get("chunk_type") or "prose")
        for c in chunks
    })

    filtered = chunks
    if chunk_type:
        filtered = [
            c for c in chunks
            if ((c.get("metadata") or {}).get("chunk_type") or "prose") == chunk_type
        ]

    page = filtered[offset : offset + limit]
    items = []
    for c in page:
        meta = c.get("metadata") or {}
        vec = c.get("vector") or []
        items.append({
            "chunk_id":       meta.get("chunk_id") or c.get("id") or "",
            "text":           c.get("text", ""),
            "token_count":    c.get("token_count"),
            "page":           meta.get("page") or c.get("page"),
            "heading_path":   meta.get("heading_path") or c.get("heading_path"),
            "chunk_type":     meta.get("chunk_type"),
            "table_name":     meta.get("table_name"),
            "doc_id":         meta.get("doc_id"),
            "vector_preview": [round(float(x), 5) for x in vec[:8]] if vec else [],
            "has_vector":     bool(vec) and len(vec) >= 100,
        })

    return {
        "job_id": job_id,
        "total": len(filtered),
        "offset": offset,
        "limit": limit,
        "chunk_types": distinct_types,
        "items": items,
    }


@app.get(
    "/api/jobs/{job_id}/inspect/sql/{table_name}",
    response_model=InspectSqlResponse,
    tags=["Jobs"],
    summary="Paginated rows from one SQLite table",
    responses={
        200: {"description": "Rows from the named table."},
        404: {"description": "Job or table not found.", "model": ErrorResponse},
    },
)
async def inspect_sql_table(
    job_id: str,
    table_name: str,
    offset: int = 0,
    limit: int = 100,
):
    """Read rows directly from the per-job SQLite database.

    Tables are named `doc_table_1`, `doc_table_2`, … — the same identifier
    used by the Qdrant `metadata.table_name` payload field and the KG's
    `Table:doc_table_N` nodes.
    """
    import sqlite3
    limit = max(1, min(limit, 1000))
    offset = max(0, offset)

    # Whitelist: only allow real tables from sql_registry to prevent injection
    import services.job_cache as cache
    registry = cache.get(job_id, "sql_registry") or cache.get(job_id, "d_sql_registry") or {}
    if table_name not in registry:
        raise HTTPException(status_code=404, detail=f"Table '{table_name}' not in this job's SQL registry")

    job_dir = Path(__file__).resolve().parent / "data" / "jobs" / job_id
    db_path = None
    for candidate in ("tables_.db", "tables.db", "tables_d_.db"):
        p = job_dir / candidate
        if p.exists():
            db_path = p
            break
    if db_path is None:
        raise HTTPException(status_code=404, detail="No SQLite database for this job")

    try:
        conn = sqlite3.connect(str(db_path))
        conn.row_factory = sqlite3.Row
        columns = [r[1] for r in conn.execute(f'PRAGMA table_info("{table_name}")').fetchall()]
        total = conn.execute(f'SELECT COUNT(*) FROM "{table_name}"').fetchone()[0]
        cur = conn.execute(f'SELECT * FROM "{table_name}" LIMIT ? OFFSET ?', (limit, offset))
        rows = []
        for i, row in enumerate(cur.fetchall()):
            rows.append({"row_index": offset + i, "cells": {k: row[k] for k in row.keys()}})
        conn.close()
    except sqlite3.Error as exc:
        raise HTTPException(status_code=500, detail=f"SQLite error: {exc}")

    return {
        "job_id": job_id,
        "table_name": table_name,
        "columns": columns,
        "total": total,
        "offset": offset,
        "limit": limit,
        "rows": rows,
    }


@app.get(
    "/api/jobs/{job_id}/inspect/kg",
    response_model=InspectKgResponse,
    tags=["Jobs"],
    summary="Paginated knowledge graph nodes (and optional node-focused detail)",
    responses={
        200: {"description": "Nodes + per-type counts. Pass `focus=<key>` for one node's neighbours."},
        404: {"description": "No knowledge graph for this job.", "model": ErrorResponse},
    },
)
async def inspect_kg(
    job_id: str,
    offset: int = 0,
    limit: int = 100,
    node_type: str | None = None,
    q: str | None = None,
    focus: str | None = None,
):
    """List every node in the per-job knowledge graph.

    Filter by `node_type` (`document`, `table`, `chunk`, or `entity`), or
    pass `q=` to substring-match the node key. Pass `focus=<exact key>` to
    also receive that node's full neighbour list + edge weights — the
    fastest path to verify "which entities are reachable from
    `Table:doc_table_3`?"
    """
    import services.job_cache as cache
    limit = max(1, min(limit, 500))
    offset = max(0, offset)

    kg = cache.get(job_id, "knowledge_graph") or cache.get(job_id, "d_knowledge_graph")
    if kg is None or not hasattr(kg, "number_of_nodes"):
        raise HTTPException(status_code=404, detail="No knowledge graph for this job")

    def _serialize_node(key: str) -> dict:
        attrs = dict(kg.nodes[key])
        return {
            "key": key,
            "type": attrs.pop("type", "unknown"),
            "label": attrs.pop("label", None),
            "text": attrs.pop("text", None) or attrs.pop("chunk_id", None) or attrs.pop("table_name", None) or attrs.pop("doc_id", None),
            "attrs": attrs,
            "degree": kg.degree(key),
        }

    # Per-type counts
    type_counts: dict[str, int] = {}
    all_keys: list[str] = []
    for key, attrs in kg.nodes(data=True):
        t = attrs.get("type", "unknown")
        type_counts[t] = type_counts.get(t, 0) + 1
        if node_type and t != node_type:
            continue
        if q and q.lower() not in key.lower():
            continue
        all_keys.append(key)

    # Stable sort: type first, then key — keeps document/table/chunk/entity grouped
    type_order = {"document": 0, "table": 1, "chunk": 2, "entity": 3}
    all_keys.sort(key=lambda k: (type_order.get(kg.nodes[k].get("type", ""), 99), k))

    page_keys = all_keys[offset : offset + limit]
    items = [_serialize_node(k) for k in page_keys]

    # Optional: detail block for one focused node
    detail = None
    if focus and kg.has_node(focus):
        neighbours = list(kg.neighbors(focus))
        edges = []
        for nb in neighbours:
            edata = kg.get_edge_data(focus, nb) or {}
            edges.append({
                "source": focus,
                "target": nb,
                "weight": edata.get("weight", 1),
                "rel": edata.get("rel"),
            })
        detail = {
            "node": _serialize_node(focus),
            "neighbours": [_serialize_node(n) for n in neighbours],
            "edges": edges,
        }

    return {
        "job_id": job_id,
        "total_nodes": kg.number_of_nodes(),
        "total_edges": kg.number_of_edges(),
        "node_types": type_counts,
        "offset": offset,
        "limit": limit,
        "items": items,
        "detail": detail,
    }


@app.get(
    "/api/jobs/{job_id}/file",
    tags=["Files"],
    summary="Download the original source file",
    responses={
        200: {"description": "Original file bytes with original MIME type, `Content-Disposition: inline`."},
        404: {"description": "File not in memory (e.g., backend restarted after upload).", "model": ErrorResponse},
    },
)
async def get_job_file(job_id: str):
    """Return the unmodified source document.

    File bytes are kept in memory for the lifetime of the backend process. If
    the backend restarts, this returns 404 — only the ingested artifacts
    persist on disk.
    """
    if job_id not in _job_files:
        raise HTTPException(status_code=404, detail="File not available")
    f = _job_files[job_id]
    return Response(
        content=f["bytes"],
        media_type=f["mime"],
        headers={"Content-Disposition": f'inline; filename="{f["filename"]}"'},
    )


@app.get(
    "/api/jobs/{job_id}/rendered",
    tags=["Files"],
    summary="Render the source file as browser-viewable HTML",
    responses={
        200: {"description": "HTML representation (for office docs) or original bytes (PDF, HTML)."},
        404: {"description": "File not available (see /api/jobs/{job_id}/file).", "model": ErrorResponse},
    },
)
async def get_job_rendered(job_id: str):
    """Return the document as HTML for in-browser preview regardless of source format.

    - PDF / HTML: returned as-is (browsers render natively)
    - DOCX / PPTX / XLSX: converted to HTML server-side via mammoth / python-docx
      / openpyxl helpers — best-effort, layout will differ from the original
    """
    if job_id not in _job_files:
        raise HTTPException(status_code=404, detail="File not available")
    f = _job_files[job_id]
    mime: str = f["mime"]
    data: bytes = f["bytes"]
    filename: str = f["filename"]

    # PDF and HTML serve as-is; browser renders them natively
    if mime == "application/pdf" or mime.startswith("text/html"):
        return Response(content=data, media_type=mime,
                        headers={"Content-Disposition": f'inline; filename="{filename}"'})

    html = await _convert_to_html(data, mime, filename)
    return Response(content=html.encode(), media_type="text/html")


_HTML_WRAPPER = """<!doctype html><html><head><meta charset="utf-8">
<style>
  body{{font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',sans-serif;
        margin:0;padding:24px 32px;background:#fff;color:#1a1a1a;line-height:1.6;font-size:14px}}
  h1{{font-size:1.5em;margin:.5em 0}}h2{{font-size:1.2em;margin:.5em 0}}
  h3{{font-size:1em;margin:.4em 0}}
  table{{border-collapse:collapse;width:100%;margin:1em 0;font-size:13px}}
  th,td{{border:1px solid #d1d5db;padding:6px 10px;text-align:left}}
  th{{background:#f3f4f6;font-weight:600}}
  tr:nth-child(even){{background:#f9fafb}}
  pre,code{{background:#f3f4f6;border-radius:4px;padding:2px 6px;font-size:12px}}
  pre{{padding:12px;overflow-x:auto}}
  .slide{{border:1px solid #e5e7eb;border-radius:8px;padding:20px;margin:16px 0;background:#fafafa}}
  .slide-num{{font-size:11px;color:#9ca3af;margin-bottom:8px}}
  blockquote{{border-left:3px solid #6366f1;padding-left:12px;color:#374151;margin:8px 0}}
</style></head><body>{body}</body></html>"""


async def _convert_to_html(data: bytes, mime: str, filename: str) -> str:
    import io

    # ── DOCX → HTML via mammoth ──────────────────────────────────────────────
    if "wordprocessingml" in mime:
        try:
            import mammoth
            result = mammoth.convert_to_html(io.BytesIO(data))
            return _HTML_WRAPPER.format(body=result.value)
        except Exception as e:
            return _HTML_WRAPPER.format(body=f"<p>DOCX conversion failed: {e}</p>")

    # ── XLSX → HTML tables via openpyxl ─────────────────────────────────────
    if "spreadsheetml" in mime:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = list(ws.iter_rows(values_only=True))
                if not rows:
                    continue
                parts.append(f"<h2>{sheet_name}</h2><table>")
                header = rows[0]
                parts.append("<tr>" + "".join(f"<th>{v or ''}</th>" for v in header) + "</tr>")
                for row in rows[1:]:
                    if any(v is not None for v in row):
                        parts.append("<tr>" + "".join(f"<td>{v if v is not None else ''}</td>" for v in row) + "</tr>")
                parts.append("</table>")
            wb.close()
            return _HTML_WRAPPER.format(body="\n".join(parts) or "<p>Empty workbook</p>")
        except Exception as e:
            return _HTML_WRAPPER.format(body=f"<p>XLSX conversion failed: {e}</p>")

    # ── PPTX → HTML slide cards via python-pptx ──────────────────────────────
    if "presentationml" in mime:
        try:
            from pptx import Presentation
            from pptx.util import Pt
            prs = Presentation(io.BytesIO(data))
            parts = [f"<h1>{filename}</h1>"]
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if not line:
                            continue
                        # Guess heading by font size or bold
                        is_heading = any(
                            (run.font.size and run.font.size >= Pt(18)) or run.font.bold
                            for run in para.runs if para.runs
                        )
                        tag = "h3" if is_heading else "p"
                        texts.append(f"<{tag}>{line}</{tag}>")
                body = "\n".join(texts) if texts else "<p><em>(no text on this slide)</em></p>"
                parts.append(f'<div class="slide"><div class="slide-num">Slide {i}</div>{body}</div>')
            return _HTML_WRAPPER.format(body="\n".join(parts))
        except Exception as e:
            return _HTML_WRAPPER.format(body=f"<p>PPTX conversion failed: {e}</p>")

    # ── Plain text fallback ───────────────────────────────────────────────────
    try:
        text = data.decode("utf-8", errors="replace")
    except Exception:
        text = "(binary file — cannot display)"
    escaped = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    return _HTML_WRAPPER.format(body=f"<pre>{escaped}</pre>")


@app.get(
    "/api/demo-docs",
    tags=["Demo Docs"],
    summary="List bundled demo documents",
    responses={200: {"description": "Catalog of demo docs with type + domain metadata."}},
)
async def list_demo_docs():
    """Return the catalog of 6 demo documents shipped with the prototype.

    Use any `filename` from this list with `POST /api/jobs` (form field
    `demo_doc=<filename>`) to ingest without uploading.

    `has_ground_truth=true` means L3 evaluation has expected answers stored
    for that doc (used when running the eval harness).
    """
    return [
        {"id": "01", "filename": "01_pharmaceutical_trial.pdf",  "doc_type": "research_paper",  "domain": "medical",    "description": "Multi-column clinical trial PDF with merged-cell tables",         "has_ground_truth": True},
        {"id": "02", "filename": "02_financial_model.xlsx",      "doc_type": "financial_report", "domain": "financial",  "description": "6-sheet SaaS financial model with cross-sheet formulas",        "has_ground_truth": True},
        {"id": "03", "filename": "03_vendor_contract.docx",      "doc_type": "contract",         "domain": "legal",      "description": "Enterprise vendor contract with H1→H3 hierarchy and SLA appendix", "has_ground_truth": True},
        {"id": "04", "filename": "04_technical_spec.html",       "doc_type": "technical_spec",   "domain": "technical",  "description": "API docs with code blocks, endpoint tables, error codes",        "has_ground_truth": True},
        {"id": "05", "filename": "05_board_presentation.pptx",   "doc_type": "presentation",     "domain": "financial",  "description": "22-slide Series B board deck with image-only financial tables",  "has_ground_truth": True},
        {"id": "06", "filename": "06_vision_ocr_demo.pdf",       "doc_type": "financial_report", "domain": "financial",  "description": "4-page PDF: typed text + table + chart image + scanned memo page", "has_ground_truth": False},
    ]


@app.get(
    "/api/demo-docs/{filename}",
    tags=["Demo Docs"],
    summary="Download a demo document",
    responses={
        200: {"description": "Raw file bytes (octet-stream)."},
        400: {"description": "Path traversal attempt rejected.", "model": ErrorResponse},
        404: {"description": "Demo doc not found.", "model": ErrorResponse},
    },
)
async def download_demo_doc(filename: str):
    """Download one of the bundled demo documents by filename.

    Use the `filename` field from `GET /api/demo-docs`. Path-traversal attempts
    (`../`, absolute paths) are rejected with HTTP 400.
    """
    path = DEMO_DOCS_DIR / filename
    if not path.exists() or not path.is_file():
        raise HTTPException(status_code=404, detail="File not found")
    # Prevent path traversal
    if path.parent.resolve() != DEMO_DOCS_DIR.resolve():
        raise HTTPException(status_code=400, detail="Invalid filename")
    return FileResponse(str(path), filename=filename, media_type="application/octet-stream")


# ── WebSocket ──────────────────────────────────────────────────────────────────

@app.websocket("/ws/{job_id}")
async def websocket_endpoint(websocket: WebSocket, job_id: str):
    """Stream per-stage progress events for a job in real time.

    **Protocol**: client connects, server immediately replays any events that
    fired before the connection (so late-connecting clients still see the full
    pipeline), then streams new events as they happen.

    **Event shape**: see the `StageEventDoc` model in OpenAPI. Each line is a
    JSON object with `stage_id`, `stage_name`, `status`, `payload`, etc.

    **Status values**: `started` (stage entering), `running` (heartbeat —
    long-running stages emit every 30s with `_heartbeat: true` in payload),
    `completed` (success, with final payload + verification), `error`
    (with error message in payload).

    **Client doesn't need to send anything** — server pushes events;
    disconnect to stop receiving them.
    """
    await websocket.accept()
    if job_id not in _ws_connections:
        _ws_connections[job_id] = []
    _ws_connections[job_id].append(websocket)
    log.info("ws_connected", job_id=job_id)

    # Replay any events emitted before this client connected
    for event in _job_events.get(job_id, []):
        try:
            await websocket.send_text(json.dumps(event))
        except Exception:
            break

    try:
        while True:
            data = await websocket.receive_text()
            if data == '{"type":"ping"}':
                await websocket.send_text('{"type":"pong"}')
    except WebSocketDisconnect:
        _ws_connections[job_id].remove(websocket)
        log.info("ws_disconnected", job_id=job_id)


# ── Internal helpers ───────────────────────────────────────────────────────────

async def _publish_event(job_id: str, event: dict):
    # Always buffer — replayed to late-connecting WebSocket clients
    if job_id in _job_events:
        _job_events[job_id].append(event)

    msg = json.dumps(event)
    conns = _ws_connections.get(job_id, [])
    dead = []
    for ws in conns:
        try:
            await ws.send_text(msg)
        except Exception:
            dead.append(ws)
    for ws in dead:
        conns.remove(ws)


async def _run_pipeline(
    job_id: str,
    pipeline: str,
    filepath: Path | None,
    source_type: str,
    tmp_path: str | None,
):
    async def publish(event: dict):
        await _publish_event(job_id, event)

    try:
        if pipeline == "custom":
            await run_custom_pipeline(job_id, filepath, source_type, publish)
        elif pipeline == "docling":
            await run_docling_pipeline(job_id, filepath, source_type, publish)
        elif pipeline == "compare":
            await asyncio.gather(
                run_custom_pipeline(job_id, filepath, source_type, publish),
                run_docling_pipeline(job_id, filepath, source_type, publish),
            )
        elif pipeline == "agent":
            await run_agent_pipeline(job_id, filepath, source_type, publish)
        else:
            await run_mock_pipeline(job_id, pipeline, publish)

        _jobs[job_id]["status"] = "completed"
        _jobs[job_id]["completed_at"] = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
    except asyncio.CancelledError:
        # User-triggered cancel (DELETE /api/jobs/{id}) or backend shutdown.
        # Surface as 'cancelled' so the UI can distinguish from a real failure.
        log.info("pipeline_cancelled", job_id=job_id)
        _jobs[job_id]["status"] = "cancelled"
        _jobs[job_id]["completed_at"] = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
        await _publish_event(job_id, {
            "job_id": job_id, "pipeline": pipeline, "stage_id": 9999,
            "stage_name": "cancelled", "status": "error",
            "timestamp_ms": time.time() * 1000, "duration_ms": 0,
            "payload": {"reason": "User requested cancel."},
        })
        raise
    except Exception as exc:
        log.error("pipeline_error", job_id=job_id, error=str(exc))
        _jobs[job_id]["status"] = "error"
        _jobs[job_id]["error"] = str(exc)
    finally:
        # Clean up task + cancel marker so DELETE on a long-finished job 404s.
        _job_tasks.pop(job_id, None)
        _cancelled_jobs.discard(job_id)
        if tmp_path:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass
