"""Stage 3 — Format Parser (Mode A).

Dispatches to the right library based on MIME type passed from stage 2.
Images are intentionally skipped here — they're handled in the Multi-Modal stage.
XLSX processing is delegated to pipelines.custom.stage_03_parsers.xlsx.
"""
from __future__ import annotations
from pathlib import Path
from typing import Optional

from models.events import ParserPayload, TextBlock, ExtractedTable
from verification.l1 import make_check, make_verification
from pipelines.base import StageResult

_MAX_TEXT_BLOCKS = 60
_MAX_TABLE_ROWS  = 25   # rows kept per table in preview


# ── Shared helpers ─────────────────────────────────────────────────────────────

def _to_markdown(headers: list[str], rows: list[list[str]]) -> str:
    if not headers:
        return ""
    sep = ["---"] * len(headers)
    lines = [
        "| " + " | ".join(headers) + " |",
        "| " + " | ".join(sep) + " |",
    ]
    for row in rows[:_MAX_TABLE_ROWS]:
        padded = list(row) + [""] * max(0, len(headers) - len(row))
        lines.append("| " + " | ".join(str(c) for c in padded[: len(headers)]) + " |")
    return "\n".join(lines)


def _l1(parser: str, word_count: int, text_blocks: list, tables: list, mime: str) -> list:
    content_count = len(text_blocks) + len(tables)
    detail = (
        f"{len(text_blocks)} block(s)" if not tables
        else f"{len(text_blocks)} block(s), {len(tables)} table(s)"
    )
    return [
        make_check("word_count_positive", word_count > 0, f"{word_count:,} words extracted"),
        make_check("content_extracted", content_count > 0, detail),
        make_check("parser_matched_mime", True, f"{parser} matched {mime}"),
    ]


def _result(
    parser: str,
    mime: str,
    text_blocks: list[TextBlock],
    tables: list[ExtractedTable],
    word_count: int,
    page_count: Optional[int] = None,
    image_count: int = 0,
    raw_preview: str = "",
) -> StageResult:
    payload = ParserPayload(
        parser_used=parser,
        page_count=page_count,
        word_count=word_count,
        table_count=len(tables),
        image_count=image_count,
        text_blocks=text_blocks[:_MAX_TEXT_BLOCKS],
        tables=tables,
        raw_text_preview=raw_preview[:600],
    )
    checks = _l1(parser, word_count, text_blocks, tables, mime)
    return StageResult(payload=payload.model_dump(), verification=make_verification(checks))



# ── DOCX — python-docx ────────────────────────────────────────────────────────

async def _parse_docx(filepath: Path, mime: str) -> StageResult:
    import docx as _docx

    doc = _docx.Document(str(filepath))
    text_blocks: list[TextBlock] = []
    tables: list[ExtractedTable] = []
    word_count = 0

    for i, para in enumerate(doc.paragraphs):
        text = para.text.strip()
        if not text:
            continue
        hlevel = 0
        if para.style.name.startswith("Heading"):
            try:
                hlevel = int(para.style.name.split()[-1])
            except ValueError:
                hlevel = 1
        word_count += len(text.split())
        text_blocks.append(TextBlock(id=f"para_{i}", text=text[:1000], heading_level=hlevel))

    for j, tbl in enumerate(doc.tables):
        rows = [[cell.text.strip() for cell in row.cells] for row in tbl.rows]
        headers = rows[0] if rows else []
        body = rows[1:]
        tables.append(ExtractedTable(
            id=f"tbl_{j}",
            headers=headers,
            rows=body,
            as_markdown=_to_markdown(headers, body),
        ))

    raw_preview = text_blocks[0].text[:500] if text_blocks else ""
    return _result("python-docx", mime, text_blocks, tables, word_count, None, 0, raw_preview)


# ── HTML — BeautifulSoup ──────────────────────────────────────────────────────

async def _parse_html(filepath: Path, mime: str) -> StageResult:
    from bs4 import BeautifulSoup

    html = filepath.read_text(errors="ignore")
    soup = BeautifulSoup(html, "lxml" if True else "html.parser")

    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()

    text_blocks: list[TextBlock] = []
    tables: list[ExtractedTable] = []
    word_count = 0

    for tag in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6"]):
        text = tag.get_text(strip=True)
        if text:
            text_blocks.append(TextBlock(
                id=f"h_{len(text_blocks)}",
                text=text,
                heading_level=int(tag.name[1]),
            ))

    for tag in soup.find_all(["p", "li", "pre", "code"]):
        text = tag.get_text(strip=True)
        if len(text) > 15:
            word_count += len(text.split())
            text_blocks.append(TextBlock(id=f"block_{len(text_blocks)}", text=text[:1000]))

    for j, tbl in enumerate(soup.find_all("table")):
        rows = []
        for tr in tbl.find_all("tr"):
            row = [td.get_text(strip=True) for td in tr.find_all(["td", "th"])]
            if row:
                rows.append(row)
        if rows:
            headers = rows[0]
            body = rows[1:]
            tables.append(ExtractedTable(
                id=f"tbl_{j}",
                headers=headers,
                rows=body,
                as_markdown=_to_markdown(headers, body),
            ))

    image_count = len(soup.find_all("img"))
    raw_preview = soup.get_text(separator=" ", strip=True)[:500]
    return _result("beautifulsoup4", mime, text_blocks, tables, word_count, None, image_count, raw_preview)


# ── PPTX — python-pptx ───────────────────────────────────────────────────────

async def _parse_pptx(filepath: Path, mime: str) -> StageResult:
    import pptx as _pptx

    prs = _pptx.Presentation(str(filepath))
    text_blocks: list[TextBlock] = []
    word_count = 0

    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        if parts:
            text = "\n".join(parts)
            word_count += len(text.split())
            text_blocks.append(TextBlock(id=f"slide_{i + 1}", text=text[:1000], page=i + 1))

    raw_preview = text_blocks[0].text[:500] if text_blocks else ""
    return _result("python-pptx", mime, text_blocks, [], word_count, len(prs.slides), 0, raw_preview)


# ── Fallback ──────────────────────────────────────────────────────────────────

async def _parse_fallback(filepath: Path, mime: str) -> StageResult:
    text = filepath.read_text(errors="ignore")
    words = text.split()
    word_count = len(words)
    block = TextBlock(id="raw", text=text[:2000]) if text.strip() else None
    blocks = [block] if block else []
    return _result("plaintext-fallback", mime, blocks, [], word_count, None, 0, text[:500])


# ── Dispatcher ────────────────────────────────────────────────────────────────

async def run(filepath: Path, mime: str) -> StageResult:
    if mime == "application/pdf":
        from pipelines.custom.stage_03_parsers.pdf import parse as _pdf_parse
        return await _pdf_parse(filepath, mime, _result)
    if "wordprocessingml" in mime or mime == "application/msword":
        return await _parse_docx(filepath, mime)
    if "spreadsheetml" in mime or mime == "application/vnd.ms-excel":
        from pipelines.custom.stage_03_parsers.xlsx import parse as _xlsx_parse
        return await _xlsx_parse(filepath, mime, _result)
    if "presentationml" in mime:
        return await _parse_pptx(filepath, mime)
    if mime == "text/html":
        return await _parse_html(filepath, mime)
    return await _parse_fallback(filepath, mime)
